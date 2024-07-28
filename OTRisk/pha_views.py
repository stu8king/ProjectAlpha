import base64
import concurrent.futures
import hashlib
import json
import logging
import math
import os
import re
import tempfile
import time
import uuid
import datetime
from decimal import Decimal
from urllib.parse import urljoin
from datetime import datetime

import openai
import requests
import urllib3
from django.contrib.auth.decorators import login_required
from django.contrib.auth.models import User
from django.core.exceptions import ObjectDoesNotExist
from django.core.serializers import serialize
from django.db import transaction
from django.db.models import Avg, Sum, F, Count, Subquery, OuterRef, Value, IntegerField, Q
from django.db.models.functions import Coalesce
from django.forms import model_to_dict
from django.http import FileResponse, HttpResponse
from django.http import JsonResponse
from django.shortcuts import render, get_object_or_404, redirect
from django.utils import timezone

from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST, require_http_methods
from pinecone import Pinecone, ServerlessSpec
from pptx import Presentation
from pptx.util import Pt
from reportlab.graphics import renderPDF
from requests.exceptions import RequestException, JSONDecodeError
from urllib3.exceptions import InsecureRequestWarning

from OTRisk.models.Model_CyberPHA import tblIndustry, tblCyberPHAHeader, tblZones, tblStandards, \
    tblCyberPHAScenario, vulnerability_analysis, tblAssetType, MitreControlAssessment, \
    SECURITY_LEVELS, ScenarioConsequences, user_scenario_audit, auditlog, CyberPHAModerators, \
    WorkflowStatus, APIKey, CyberPHA_Group, ScenarioBuilder, PHA_Safeguard, CyberSecurityInvestment, UserScenarioHash, \
    CyberPHACybersecurityDefaults, PHA_Observations, Country, OTVendor, Facility, FacilityType
from OTRisk.models.model_assessment import SelfAssessment, AssessmentQuestion
from OTRisk.models.raw import MitreICSMitigations, RAActions
from OTRisk.models.raw import SecurityControls
from ProjectAlpha import settings
from ProjectAlpha.settings import BASE_DIR
from accounts.models import UserProfile, Organization
from accounts.views import get_client_ip
from .dashboard_views import get_user_organization_id, get_organization_users
from .forms import VulnerabilityAnalysisForm
from OTRisk.models.darktraceapi import DarktraceAPI
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.shapes import Drawing, Rect, Image, Line, Circle
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import letter
from reportlab.lib.sequencer import getSequencer
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch, mm
from reportlab.graphics.charts.textlabels import Label
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen.canvas import Canvas
from reportlab.platypus import SimpleDocTemplate, PageBreak, Image, Frame, Image, Table, TableStyle, BaseDocTemplate, \
    PageTemplate, ListFlowable, ListItem, Flowable
from reportlab.platypus import Table, TableStyle, Paragraph, Spacer
from io import BytesIO
from reportlab.graphics.shapes import Drawing, String
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.widgetbase import Widget
from diagrams import Diagram, Edge, Node
from diagrams.custom import Custom


def get_api_key(service_name):
    try:
        key_record = APIKey.objects.get(service_name=service_name)
        return key_record.key
    except ObjectDoesNotExist:
        # Handle the case where the key is not found
        return None


openai.api_key = get_api_key('openai')


def validate_and_format_date(date_str, default_date='2001-01-01', date_format='%Y-%m-%d'):
    """
    Validates and formats a date string into a timezone-aware datetime object.
    If the date string is empty, the function returns the current date.

    :param date_str: The date string to validate and format.
    :param default_date: The default date to return if date_str is invalid or empty.
    :param date_format: The format to which the date string should be formatted.
    :return: A timezone-aware datetime object representing the validated and formatted date.
    """
    if date_str:
        try:
            # Attempt to parse the date string using the specified format
            valid_date = datetime.datetime.strptime(date_str, date_format)
        except ValueError:
            # If parsing fails, use the default date
            valid_date = datetime.datetime.strptime(default_date, date_format)
    else:
        # If the date string is empty, use the current date
        valid_date = datetime.datetime.now()

    # Make the datetime object timezone-aware
    timezone_aware_date = timezone.make_aware(valid_date, timezone.get_default_timezone())
    return timezone_aware_date


@login_required
def iotaphamanager(request, record_id=None):
    auto_start = request.GET.get('autoStart', False)

    pha_header = None
    new_record_id = None  # Initialize new_record_id to None
    ### annual_revenue_str = "$0"
    ### annual_revenue_str = "$0" coho_str = "$0"

    if request.method == 'POST':
        is_new_record = False  # Initialize flag
        title = request.POST.get('txtTitle')
        ### facility_name = request.POST.get('txtFacility')
        # Check for duplicate record

        pha_id = request.POST.get('txtHdnCyberPHAID')
        if pha_id and int(pha_id) > 0:
            # Update existing record
            pha_header, created = tblCyberPHAHeader.objects.get_or_create(ID=pha_id)
        else:
            ### duplicate_record = tblCyberPHAHeader.objects.filter(title=title, FacilityName=facility_name).exists()
            ### if duplicate_record:
            ###    return redirect('OTRisk:iotaphamanager')
            # Create a new record
            is_new_record = True
            pha_header = tblCyberPHAHeader()

        pha_header.facility_id = int(request.POST.get('facilityid'))
        PHATitle = request.POST.get('txtTitle')
        pha_header.title = PHATitle if PHATitle else "Not Given"
        PHALeaderName = request.POST.get('txtLeader')
        pha_header.PHALeader = PHALeaderName if PHALeaderName else "Not Given"
        PHALeaderEmail = request.POST.get('txtLeaderEmail')
        pha_header.PHALeaderEmail = PHALeaderEmail if PHALeaderEmail else "Not Given"
        pha_header.FacilityName = ''  ### request.POST.get('txtFacility')
        pha_header.Industry = ''  ### request.POST.get('selIndustry')
        pha_header.FacilityType = ''  ### request.POST.get('selFacilityType')

        pha_header.AssessmentUnit = request.POST.get('txtUnit')

        selZone_value = request.POST.get('selZone')
        pha_header.AssessmentZone = selZone_value if selZone_value else "None"

        start_date_str = request.POST.get('txtStartDate')
        pha_header.AssessmentStartDate = validate_and_format_date(start_date_str)

        # Validate and format AssessmentEndDate
        end_date_str = request.POST.get('txtEndDate')
        pha_header.AssessmentEndDate = validate_and_format_date(end_date_str)

        pha_header.facilityAddress = ''  ### request.POST.get('txtAddress')

        if is_new_record:
            risk_profile_data = facility_risk_profile_newrecord(request.user.id,
                                                                int(request.POST.get('facilityid')),

                                                                int(request.POST.get('assessment_id') or 0),
                                                                int(request.POST.get('sl') or 0),
                                                                request.POST.get('ir_plan'),
                                                                request.POST.get('ir_plan_ut'),
                                                                request.POST.get('ir_plan_date')
                                                                )
            ### pha_header.safetySummary = risk_profile_data['safety_summary']
            ### pha_header.chemicalSummary = risk_profile_data['chemical_summary']
            ### pha_header.physicalSummary = risk_profile_data['physical_security_summary']
            pha_header.otherSummary = risk_profile_data['other_summary']
            pha_header.threatSummary = risk_profile_data['threatSummary']
            pha_header.insightSummary = risk_profile_data['insightSummary']
            pha_header.strategySummary = risk_profile_data['strategySummary']
            ### pha_header.complianceSummary = risk_profile_data['compliance_summary']
            ### pha_header.pha_score = risk_profile_data['pha_score']

        else:
            pha_header.safetySummary = ''  ###  request.POST.get('txtSafetySummary')
            pha_header.chemicalSummary = ''  ### request.POST.get('txtChemical')
            pha_header.physicalSummary = ''  ### request.POST.get('txtPhysical')
            pha_header.otherSummary = request.POST.get('txtOther')
            pha_header.threatSummary = request.POST.get('threatSummary')
            pha_header.insightSummary = request.POST.get('insightSummary')
            pha_header.strategySummary = request.POST.get('strategySummary')
            pha_header.complianceSummary = ''  ### request.POST.get('txtCompliance
            pha_header.pha_score = 0

            pha_header.country = ""
            validated_date = validate_and_format_date(start_date_str)
            pha_header.Date = validated_date if validated_date else datetime.now()
            pha_header.EmployeesOnSite = 0  ### int(request.POST.get('txtEmployees') or 0)
            pha_header.facilityAQI = 0  ### request.POST.get('txthdnAQI')
            pha_header.facilityCity = ""
            pha_header.facilityCode = ""
            pha_header.facilityLat = ""

            pha_header.facilityLong = ""
            pha_header.facilityState = ""
            pha_header.shift_model = ""  #### request.POST.get('shift_model')
        try:
            assessment_id = int(request.POST.get('assessment_id')) if request.POST.get('assessment_id') else None
        except ValueError:
            assessment_id = None
        pha_header.assessment = assessment_id
        pha_header.last_assessment_score = int(request.POST.get('last_assessment_score') or 0)

        pha_header.last_assessment_summary = request.POST.get('last_assessment_summary') or ''
        pha_header.npm = 0

        # Continue with the rest of the processing

        pha_header.sl_t = request.POST.get('selSL')
        pha_header.FacilityID = 0
        pha_header.Deleted = 0

        annual_revenue_str = 0  ### request.POST.get('annual_revenue', '')
        coho_str = ''  ### request.POST.get('coho', '')
        # Strip out $ and , characters
        cleaned_annual_revenue_str = ''  ### .join(filter(str.isdigit, annual_revenue_str))
        cleaned_coho_str = ''  ### .join(filter(str.isdigit, coho_str))

        # Convert the cleaned string to an integer
        annual_revenue_int = 0  # Or handle this situation differently if needed
        coho_int = 0  # Or handle this situation differently if needed
        # Save to your model
        pha_header.annual_revenue = 0  ### annual_revenue_int
        pha_header.coho = coho_int

        cyber_insurance_value = request.POST.get('cyber_insurance')
        pha_header.cyber_insurance = False if cyber_insurance_value is None else bool(cyber_insurance_value)
        pha_header.is_default = request.POST.get('defaultFacility') == 'on'
        has_ir_plan = request.POST.get('ir_plan') == 'on'  # Checkbox 'on' if checked
        ir_plan_date_str = request.POST.get('ir_plan_date')
        ir_plan_never_tested = request.POST.get('ir_plan_ut') == 'on'
        pha_header.exalens_ip = request.POST.get('exalensIpAddress')
        pha_header.exalens_api = request.POST.get('exalensApiKey')
        pha_header.exalens_client = request.POST.get('exalensClientId')
        pha_header.exalens_status = request.POST.get('hdn_exalens_status')
        pha_header.exalens_risk = request.POST.get('hdn_exalens_risk')
        exalens_risk_score_string = request.POST.get('hdn_exalens_score')
        try:
            exalens_risk_score_int = int(exalens_risk_score_string)
        except ValueError:  # Handle cases where the input might still not be a valid integer
            exalens_risk_score_int = 0
        pha_header.exalens_score = exalens_risk_score_int

        pha_header.darktrace_client = request.POST.get('darktrace_host')
        pha_header.darktrace_public_api = request.POST.get('darktrace_public_key')
        pha_header.darktrace_private_api = request.POST.get('darktrace_private_key')
        pha_header.darktrace_status = request.POST.get('hdn_darktrace_status')
        pha_header.darktrace_risk = request.POST.get('hdn_darktrace_risk')
        pha_header.darktrace_tactics_report = request.POST.get('hdn_darktrace_tactics')
        pha_header.darktrace_mitre_report = request.POST.get('hdn_darktrace_mitre')

        darktrace_risk_score_string = request.POST.get('hdn_darktrace_score')
        try:
            darktrace_risk_score_int = int(darktrace_risk_score_string)
        except ValueError:  # Handle cases where the input might still not be a valid integer
            darktrace_risk_score_int = 0
        pha_header.darktrace_score = darktrace_risk_score_int

        # Convert ir_plan_date from string to date object, handle empty string
        if ir_plan_date_str:
            ir_plan_date = timezone.datetime.strptime(ir_plan_date_str, '%Y-%m-%d').date()
            pha_header.plan_last_tested_date = ir_plan_date
            pha_header.plan_never_tested = False  # If a date is provided, plan has been tested
        else:
            # If no date provided and plan exists, consider if it's marked as never tested
            if has_ir_plan:
                pha_header.plan_never_tested = ir_plan_never_tested
            else:
                # If no IR plan, reset related fields
                pha_header.plan_last_tested_date = None
                pha_header.plan_never_tested = True

        pha_header.has_incident_response_plan = has_ir_plan
        pha_header.UserID = request.user.id
        pha_header.save()

        saved_record_id = pha_header.ID
        if is_new_record:
            user_action = f"Created a new CyberPHA titled {pha_header.title}"
        else:
            user_action = f"Amended CyberPHA titled {pha_header.title}"

        cyber_pha_instance = tblCyberPHAHeader.objects.get(ID=saved_record_id)
        write_to_audit(
            user_id=request.user,
            user_action=user_action,
            user_ip=get_client_ip(request),
            cyberPHAID=cyber_pha_instance,
            cyberPHAScenario=None,
            qraw=None
        )

        #### Save investment information

        if pha_id and int(pha_id) > 0:
            CyberSecurityInvestment.objects.filter(cyber_pha_header=pha_header).delete()

        # Extract and process investment information from POST data

        vendor_names = request.POST.getlist('vendor_name[]')
        product_names = request.POST.getlist('product_name[]')

        for vendor, product in zip(vendor_names, product_names):
            vendor = vendor if vendor else ""
            product = product if product else ""

            CyberSecurityInvestment.objects.create(
                cyber_pha_header=pha_header,
                vendor_name=vendor,
                product_name=product
            )

        #### End save investments

        #### End save risk tolerance data
        # After saving pha_header, handle CyberPHACybersecurityDefaults
        if request.method == 'POST':
            # Assuming pha_header is the instance of tblCyberPHAHeader you've just created or updated
            defaults_data = {
                'overall_aggregate_limit': request.POST.get('overall_aggregate_limit'),
                'per_claim_limit': request.POST.get('per_claim_limit'),
                'deductible_amount': request.POST.get('deductible_amount'),
                'first_party_coverage': request.POST.get('first_party_coverage') == 'on',
                'third_party_coverage': request.POST.get('third_party_coverage') == 'on',
                'security_event_liability': request.POST.get('security_event_liability') == 'on',
                'privacy_regulatory_actions': request.POST.get('privacy_regulatory_actions') == 'on',
                'cyber_extortion_coverage': request.POST.get('cyber_extortion_coverage') == 'on',
                'data_breach_response_coverage': request.POST.get('data_breach_response_coverage') == 'on',
                'business_interruption_coverage': request.POST.get('business_interruption_coverage') == 'on',
                'dependent_business_coverage': request.POST.get('dependent_business_coverage') == 'on',
                'data_recovery': request.POST.get('data_recovery') == 'on',
                'hardware_replacement': request.POST.get('hardware_replacement') == 'on',
                'reputation_harm': request.POST.get('reputation_harm') == 'on',
                'media_liability': request.POST.get('media_liability') == 'on',
                'pci_dss': request.POST.get('pci_dss') == 'on',
                'premium_base': request.POST.get('premium_base'),
                'notification_period_days': request.POST.get('notification_period_days'),
                'cancellation_terms_days': request.POST.get('cancellation_terms_days'),
            }

            # Convert string values to appropriate types
            for key in ['overall_aggregate_limit', 'per_claim_limit', 'deductible_amount', 'premium_base']:
                defaults_data[key] = float(defaults_data[key]) if defaults_data[key] else 0.0
            for key in ['notification_period_days', 'cancellation_terms_days']:
                defaults_data[key] = int(defaults_data[key]) if defaults_data[key] else 0

            # Update or create CyberPHACybersecurityDefaults instance
            defaults_instance, created = CyberPHACybersecurityDefaults.objects.update_or_create(
                cyber_pha=pha_header,
                defaults=defaults_data
            )

        if is_new_record:
            pha_header.set_workflow_status('Started')
        # Save Workflow Status
        selected_workflow_status = request.POST.get('workflow_selector')
        if selected_workflow_status:
            pha_header.set_workflow_status(selected_workflow_status)
        # First, remove existing moderators for this PHA record
        CyberPHAModerators.objects.filter(pha_header=pha_header).delete()
        # Get the list of selected moderator IDs from the form

        selected_moderators_ids = request.POST.getlist('moderator')
        # Get the target date from the form

        target_date_str = request.POST.get('targetDate', None)
        if target_date_str:
            try:
                # Create a naive datetime object from the string
                naive_datetime = datetime.strptime(target_date_str, '%m/%d/%Y')
                # Make it timezone aware
                target_date = timezone.make_aware(naive_datetime, timezone.get_default_timezone())
            except ValueError:
                target_date = None
        else:
            target_date = None

        # Create new associations
        for moderator_id in selected_moderators_ids:
            moderator = User.objects.get(id=moderator_id)
            CyberPHAModerators.objects.create(pha_header=pha_header, moderator=moderator, target_date=target_date)

        new_record_id = pha_header.ID

    organization_id_from_session = request.session.get('user_organization')

    users_in_organization = User.objects.filter(userprofile__organization__id=organization_id_from_session)

    ra_actions_subquery = RAActions.objects.filter(phaID=OuterRef('ID')).values('phaID').annotate(
        action_count=Count('ID')).values('action_count')

    # Subquery to get the latest workflow status for each tblCyberPHAHeader record
    latest_status_subquery = WorkflowStatus.objects.filter(
        cyber_pha_header=OuterRef('pk')
    ).order_by('-timestamp').values('status')[:1]

    # Annotate the pha_header_records queryset with the latest workflow status
    pha_header_records = tblCyberPHAHeader.objects.filter(
        UserID__in=users_in_organization,
        Deleted=0
    ).annotate(
        scenario_count=Count('tblcyberphascenario'),
        ra_action_count=Coalesce(Subquery(ra_actions_subquery, output_field=IntegerField()), Value(0)),
        latest_workflow_status=Subquery(latest_status_subquery)
    )

    if record_id is not None:
        pha_header = get_object_or_404(tblCyberPHAHeader, pk=record_id)
        new_record_id = record_id
        first_record_id = record_id
    else:
        first_record = pha_header_records.first()
        first_record_id = first_record.ID if first_record else 0
        pha_header = None
    # get the list of assessments
    # 0 means a global assessment that's built in for all customer
    # any other value is a customer specific assessment so should only be visible for the customer where their id matches
    current_user_profile = UserProfile.objects.get(user=request.user)
    user_organization_id = request.session.get('user_organization', 0)  # Default to 0 if not in session
    assessments = SelfAssessment.objects.filter(
        Q(organization_id=user_organization_id)
    )

    walkdowns = SelfAssessment.objects.filter(
        organization_id=user_organization_id,
        framework__name__icontains='Walkdown'
    )

    industries = tblIndustry.objects.all().order_by('Industry')
    countries = Country.objects.all().order_by('country')
    facilities = FacilityType.objects.all().order_by('FacilityType')
    zones = tblZones.objects.all().order_by('PlantZone')
    standardslist = tblStandards.objects.all().order_by('standard')
    vendors_products = list(OTVendor.objects.all().values('vendor', 'product'))

    # Create a set for unique vendors
    unique_vendors = {vp['vendor'] for vp in vendors_products}

    # Convert the vendors and products to JSON
    vendors_json = json.dumps(vendors_products)
    unique_vendors_json = json.dumps(list(unique_vendors))
    mitigations = MitreICSMitigations.objects.all()
    anychart_key = get_api_key('anychart')
    moderators_in_organization = UserProfile.objects.filter(
        organization_id=user_organization_id,
        role_moderator=True
    )

    # Retrieve current workflow status for the pha_header
    current_workflow_status = "Started"  # Default status

    if pha_header:
        latest_status = pha_header.workflow_statuses.last()
        if latest_status:
            current_workflow_status = latest_status.status

    return render(request, 'iotaphamanager.html', {
        'pha_header_records': pha_header_records,
        'industries': industries,
        'facilities': facilities,
        'zones': zones,
        'standardslist': standardslist,
        'current_pha_header': pha_header,
        'new_record_id': new_record_id,
        'mitigations': mitigations,
        'SHIFT_MODELS': tblCyberPHAHeader.SHIFT_MODELS,
        'annual_revenue_str': '0',  # annual_revenue_str,
        'coho_str': '0',  # coho_str,
        'selected_record_id': first_record_id,
        'SECURITY_LEVELS': SECURITY_LEVELS,
        'assessments': assessments,
        'walkdowns': walkdowns,
        'moderators': moderators_in_organization,
        'current_workflow_status': current_workflow_status,
        'workflow_status_choices': WorkflowStatus.STATUS_CHOICES,
        'anychart_key': anychart_key,
        'group_types': CyberPHA_Group.GROUP_TYPES,
        'saved_record_id': new_record_id,
        'countries': countries,
        'vendors_json': vendors_json,
        'unique_vendors_json': unique_vendors_json,
        'auto_start': auto_start
    })


def get_headerrecord(request):
    record_id = request.GET.get('record_id')
    headerrecord = get_object_or_404(tblCyberPHAHeader, ID=record_id)
    # Retrieve the latest workflow status for this header record
    latest_status = headerrecord.workflow_statuses.last()
    current_workflow_status = latest_status.status if latest_status else "Started"
    # Get current group assignments
    current_groups = headerrecord.groups.all()
    current_groups_data = serialize('json', current_groups)

    # Serialize all_groups
    all_groups = CyberPHA_Group.objects.all()
    all_groups_data = serialize('json', all_groups)
    # Retrieve the CyberPHARiskTolerance record associated with this header record
    facility = get_object_or_404(Facility, id=headerrecord.facility_id)
    industry = facility.industry.Industry
    facility_type = facility.type.FacilityType
    # create a dictionary with the record data
    headerrecord_data = {
        'title': headerrecord.title,
        'facility': facility.name,
        'business_name': facility.business_name,
        'leader': headerrecord.PHALeader,
        'leaderemail': headerrecord.PHALeaderEmail,
        'Industry': industry,
        'facilitytype': facility_type,
        'unit': headerrecord.AssessmentUnit,
        'zone': headerrecord.AssessmentZone,
        'startdate': headerrecord.AssessmentStartDate.strftime('%Y-%m-%d'),
        'enddate': headerrecord.AssessmentEndDate.strftime('%Y-%m-%d'),
        'address': facility.address,
        'safetysummary': facility.safetySummary,
        'chemicalsummary': facility.type.chemical_profile,
        'physicalsummary': facility.physicalSummary,
        'othersummary': headerrecord.otherSummary,
        'compliancesummary': facility.complianceSummary,
        'threatSummary': headerrecord.threatSummary,
        'insightSummary': headerrecord.insightSummary,
        'strategySummary': headerrecord.strategySummary,
        # 'country': headerrecord.country,
        'shift_model': facility.shift_model,
        'EmployeesOnSite': facility.employees,
        'cyber_insurance': headerrecord.cyber_insurance,
        'annual_revenue': facility.revenue,
        'pha_score': facility.pha_score,
        'sl_t': headerrecord.sl_t,
        'assessment_id': headerrecord.assessment,
        'last_assessment_score': headerrecord.last_assessment_score,
        'last_assessment_summary': headerrecord.last_assessment_summary,
        'coho': facility.operating_cost,
        'npm': facility.profit_margin,
        'current_workflow_status': current_workflow_status,
        'current_groups': current_groups_data,
        'all_groups': all_groups_data,
        'group_types': CyberPHA_Group.GROUP_TYPES,
        'facilityAQI': facility.aqi_score,
        # 'facilityCity': headerrecord.facilityCity,
        # 'facilityCode': headerrecord.facilityCode,
        'facilityLat': facility.lat,
        'facilityLong': facility.lon,
        'facilityState': headerrecord.facilityState,
        'has_incident_response_plan': headerrecord.has_incident_response_plan,
        'plan_last_tested_date': headerrecord.plan_last_tested_date.strftime(
            '%Y-%m-%d') if headerrecord.plan_last_tested_date else None,
        'plan_never_tested': headerrecord.plan_never_tested,
        'is_default': headerrecord.is_default,
        'exalens_api': headerrecord.exalens_api,
        'exalens_client': headerrecord.exalens_client,
        'exalens_ip': headerrecord.exalens_ip,
        'exalens_risk': headerrecord.exalens_risk,
        'exalens_score': headerrecord.exalens_score,
        'exalens_status': headerrecord.exalens_status,
        'darktrace_public_api': headerrecord.darktrace_public_api,
        'darktrace_private_api': headerrecord.darktrace_private_api,
        'darktrace_client': headerrecord.darktrace_client,
        'darktrace_risk': headerrecord.darktrace_risk,
        'darktrace_score': headerrecord.darktrace_score,
        'darktrace_status': headerrecord.darktrace_status,
        'facility_id': headerrecord.facility_id
    }

    # Safely load JSON data for darktrace_mitre_report and darktrace_tactics_report
    def load_json_field(json_field):
        try:
            return json.loads(json_field) if json_field else ""
        except (JSONDecodeError, TypeError):
            return ""

    darktrace_mitre_report = headerrecord.darktrace_mitre_report
    darktrace_tactics_report = headerrecord.darktrace_tactics_report

    if darktrace_mitre_report is not None and darktrace_mitre_report != "":
        headerrecord_data['darktrace_mitre_report'] = load_json_field(darktrace_mitre_report)
    else:
        headerrecord_data['darktrace_mitre_report'] = None  # Or handle the absence of data as needed

    if darktrace_tactics_report is not None and darktrace_tactics_report != "":
        headerrecord_data['darktrace_tactics_report'] = load_json_field(darktrace_tactics_report)
    else:
        headerrecord_data['darktrace_tactics_report'] = None

        # Query for moderators associated with this header record
    moderators = CyberPHAModerators.objects.filter(pha_header=headerrecord)
    moderator_ids = [moderator.moderator.id for moderator in moderators]

    moderators_data = [
        {'id': moderator.moderator.id, 'name': f"{moderator.moderator.first_name} {moderator.moderator.last_name}"}
        for moderator in moderators
    ]

    # Retrieve all users from the current organization who are moderators
    organization_moderators = User.objects.filter(
        userprofile__organization_id=get_user_organization_id(request),
        userprofile__role_moderator=True
    )
    organization_moderators_data = [
        {'id': user.id, 'name': f"{user.first_name} {user.last_name}"}
        for user in organization_moderators
    ]
    cyber_pha_instance = tblCyberPHAHeader.objects.get(ID=record_id)
    # Log the user activity
    write_to_audit(
        user_id=request.user,
        user_action=f'Viewed cyberPHA: {headerrecord.title}',
        user_ip=get_client_ip(request),
        cyberPHAID=cyber_pha_instance
    )

    control_assessments = MitreControlAssessment.objects.filter(cyberPHA=headerrecord)
    # control_effectiveness = math.ceil(calculate_effectiveness(record_id))
    try:
        control_effectiveness = SelfAssessment.objects.get(id=headerrecord.assessment).score_effective
        # assessment_summary_result = assessment_summary(headerrecord.assessment)
    except SelfAssessment.DoesNotExist:
        control_effectiveness = 0
    # Create a list of dictionaries for control assessments
    control_assessments_data = []
    for assessment in control_assessments:
        assessment_data = {
            'mitigation_id': assessment.control.id,  # Assuming control has an ID field
            'effectiveness_percentage': assessment.effectiveness_percentage,
            'weighting': assessment.weighting,
            # Add other fields if needed
        }
        control_assessments_data.append(assessment_data)

    investments = CyberSecurityInvestment.objects.filter(facility=facility).values(
        'id', 'investment_type', 'vendor_name', 'product_name', 'cost', 'date'
    )
    investments_data = list(investments)

    print(headerrecord_data)

    response_data = {
        'headerrecord': headerrecord_data,
        'control_assessments': control_assessments_data,
        'control_effectiveness': control_effectiveness,
        'organization_moderators': organization_moderators_data,  # All moderators in the organization
        'current_moderators': moderators_data,  # Moderators for the specific header record
        'moderator_ids': moderator_ids,  # IDs of Moderators for the specific header record
        'investments': investments_data
    }

    return JsonResponse(response_data)


def extract_section(text, title):
    """
    Extracts a section from the provided text using the given title.
    """
    start = text.find(title)
    if start == -1:
        return ''
    start += len(title)
    end = text.find('**', start)
    return text[start:end].strip() if end != -1 else text[start:].strip()


def process_section(section_text):
    # this version of process section works for gpt-4-turbo
    """
    Processes the extracted section text to handle detailed structure,
    including sub-points and bolded text.
    """
    # Initialize an empty list to hold the processed points
    processed_points = []
    current_point_lines = []  # Temporary storage for accumulating lines of the current point

    # Split the section text into lines for processing
    lines = section_text.split('\n')

    for line in lines:
        # Check if the line starts a new main point (e.g., "1. ")
        if line.strip().startswith("1. ") or line.strip().startswith("2. ") or line.strip().startswith(
                "3. ") or line.strip().startswith("4. ") or line.strip().startswith("5. ") or line.strip().startswith(
            "6. ") or line.strip().startswith("7. ") or line.strip().startswith("8. ") or line.strip().startswith(
            "9. ") or line.strip().startswith("10. "):
            # If there's content in current_point_lines, it means the previous point is complete
            if current_point_lines:
                # Join the accumulated lines into a single string and add to the list
                processed_points.append('\n'.join(current_point_lines).strip())
                current_point_lines = [line]  # Start a new point
            else:
                # If current_point_lines is empty, this is the first point
                current_point_lines.append(line)
        else:
            # If the line does not start a new main point, it's a continuation or sub-point
            current_point_lines.append(line)

    # After the loop, add the last accumulated point to the list
    if current_point_lines:
        processed_points.append('\n'.join(current_point_lines).strip())

    # Join the processed points into a single string with double line breaks between points
    return '\n\n'.join(processed_points)


def process_section_gpt4(section_text):
    # this version of process section works for GPT4
    bullet_points = section_text.split('\n')

    # Add a line space between each bullet point
    processed_text = '\n\n'.join(bullet_point.strip() for bullet_point in bullet_points if bullet_point.strip())

    return processed_text


def make_request_with_backoff(openai_function, *args, **kwargs):
    max_attempts = 5
    base_delay = 1.0  # Base delay in seconds
    for attempt in range(max_attempts):
        try:
            return openai_function(*args, **kwargs)
        except openai.error.ServiceUnavailableError:
            sleep_time = base_delay * (2 ** attempt)

            time.sleep(sleep_time)
    raise Exception("Failed to make request after several attempts.")


def facility_threat_profile(security_level, facility, facility_type, country, industry,

                            other_summary, investment_statement,
                            has_ir_plan_str, ir_plan_never_tested_str, ir_plan_date_str, connector_risk,
                            connector_status, connector_score):
    openai_api_key = get_api_key('openai')
    openai_api_key = get_api_key('openai')
    ai_model = get_api_key('OpenAI_Model')

    # Constructing the detailed context
    context = f"""
        Analyze the cybersecurity posture of {facility}, a {facility_type} in {country}, focusing on OT cybersecurity risk mitigation. This facility is notable in the {industry} industry and has specific challenges and assets:
        OT Devices: {other_summary}
        Incident Response Plan: {has_ir_plan_str} (Last tested: {ir_plan_date_str})
        Cybersecurity Investments: {investment_statement}
        Target Security Level (SL-T) as per IEC62443-3-2: {security_level}.
        """
    connector_score = int(connector_score)
    if connector_score > 0:
        context = context + f""" 
            The facility has invested in and deployed an OT Cybersecurity advanced threat detection and monitoring solution - Exalens - to monitor the OT network. The assessment of risk based on output from Exalens is: {connector_risk} and a risk score of {connector_score}/100. The risk status is described as {connector_status} 
       """
    prompt = f"""
        {context}
     Based on the facility's profile, and investments, provide an executive-level cybersecurity analysis specifically for OT/ICS environments. The analysis should be divided into three sections:
        
        1. Cybersecurity Threats: Make an estimate of up to 10 main cybersecurity threats and the actors likely to target this facility, considering the operational technology used, country, and industry specifics. EXTRA INSTRUCTION append a probability (as a percentage) of each threat occurring in the next 12 months..
        2. Predictive Insights: Offer up to 10 insights on potential future cybersecurity events based on current data and trends.
        3. Proactive Defense Strategies: Suggest up to 10 strategies to improve the facility's cybersecurity posture and achieve the target security level.
        
        INSTRUCTION: Utilize relevant and credible sources of information and industry reports such as from Dragos, Gartner, Deloitte.
        
        Example Format:
        Section 1: Threats
        1. Threat description. (Probability: X%)
        2. Threat description. (Probability: X%)
        
        Section 2: Insights
        1. Example insight.
        2. Another example insight.
        
        Section 3: Strategies
        1. Example strategy.
        2. Another example strategy.
        
        Please follow this format for your response, without using '###', '**', or any other special formatting characters.
        """

    # API call using chat model endpoint with the correct 'messages' property
    response = make_request_with_backoff(
        openai.ChatCompletion.create,
        model=ai_model,  # Ensure to use GPT-4 model
        messages=[
            {"role": "system",
             "content": "You are a model trained to provide concise and informative responses in a specific format."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.4,
        max_tokens=4000,
    )

    full_response = response['choices'][0]['message']['content']

    # Function to extract sections from the full_response
    def extract_section1(full_response1, section_title):
        start_index = full_response1.find(section_title)
        if start_index == -1:
            return ""
        start_index += len(section_title)
        end_index = full_response1.find("Section", start_index)
        section_text = full_response1[start_index:end_index].strip()
        return section_text

    # Extracting each section based on titles
    threat_summary = extract_section1(full_response, "Threats")
    insight_summary = extract_section1(full_response, "Insights")
    strategy_summary = extract_section1(full_response, "Strategies")

    return threat_summary, insight_summary, strategy_summary


@login_required()
def facility_risk_profile(request):
    if request.method == 'GET':
        # Gather the necessary data for the risk assessment (impact scores and scenario information)
        facility_id = int(request.GET.get('facilityid'))
        facility = get_object_or_404(Facility, id=facility_id)

        Industry = facility.industry
        facility_type = facility.type
        address = facility.address
        address_parts = address.split(',')
        country = address_parts[-1].strip() if address_parts else ''

        facility_name = facility.business_name
        employees = facility.employees

        shift_model = facility.shift_model
        assessment_id = request.GET.get('assessment_id')
        investments_data = request.GET.get('investments')

        has_ir_plan = request.GET.get('has_ir_plan', 'false') == 'true'
        ir_plan_never_tested = request.GET.get('ir_plan_never_tested', 'false') == 'true'
        ir_plan_date_str = request.GET.get('ir_plan_date')
        connector_status = request.GET.get('connector_status')
        connector_risk = request.GET.get('connector_risk')
        connector_score = request.GET.get('connector_score')

        aqi = request.GET.get('aqi')
        sl = request.GET.get('sl')
        if investments_data:
            investments = json.loads(investments_data)
        else:
            investments = []

        has_ir_plan_str = 'True' if has_ir_plan else 'False'
        ir_plan_never_tested_str = 'True' if ir_plan_never_tested else 'False'

        # Generate the text statement for investments
        investment_statement = "Investments have been made in:\n"

        for idx, investment in enumerate(investments, start=1):
            vendor_name = investment.get('vendor_name', 'Unknown Vendor')
            product_name = investment.get('product_name', 'Unknown Product')
            investment_statement += f"{idx}: Vendor:{vendor_name}, Product:{product_name}.\n"

        language = request.session.get('organization_defaults', {}).get('language', 'en')  # 'en' is the default value

        # Log the user activity
        write_to_audit(
            user_id=request.user,
            user_action=f'Generated a cyberPHA risk profile for: {facility}',
            user_ip=get_client_ip(request)
        )
        # Check if Industry or facility_type are empty or None
        if not Industry or not facility_type:
            error_msg = "Missing industry or facility type. Complete all fields to get an accurate assessment"
            return JsonResponse({
                'safety_summary': error_msg,
                'chemical_summary': error_msg,
                'physical_security_summary': error_msg,
                'other_summary': error_msg
            })

        openai_api_key = get_api_key('openai')
        openai_model = get_api_key('OpenAI_Model')
        # openai_api_key = os.environ.get('OPENAI_API_KEY')
        openai.api_key = openai_api_key

        OT_query = f"CyberPHA, OT Cybersecurity, Industrial Cyber, OT Device, Industrial control systems, Manufacturing, Industry, ICS Cybersecurity"
        retrieved_chunks = query_index(OT_query)
        summarized_chunks = get_summarized_chunks(retrieved_chunks)
        OT_context = "\n\n".join(summarized_chunks)

        context = f"You are an industrial safety and hazard expert. For the {facility} {facility_type} at {address}, {country} in the {Industry} industry, with {employees} employees working a {shift_model} shift model. The local Air Quality Index is {aqi}.  "

        prompts = [
            f"INSTRUCTION: DO NOT PRINT ** characters. If any words are emphasized with **, replace them with normal text without ** characters.   indexed pinecone content for context: {OT_context}.{context}, list of OT devices expected to be operating on the industrial networks at the given facility. Use the pinecone index for more context. .\n\nExample Format:\n 1. OT device (brief and concise purpose of device).\n 2. Another OT device (brief and concise purpose of device).",
        ]

        responses = []

        # Loop through the prompts and make an API call for each one
        def fetch_response(prompt):
            return openai.ChatCompletion.create(
                # model="gpt-4",
                model=openai_model,
                messages=[
                    {"role": "system",
                     "content": "You are a model trained to provide concise responses. Please provide a concise numbered bullet-point list based on the given statement."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_tokens=4000

            )

            # Use ThreadPoolExecutor to parallelize the API calls

        with concurrent.futures.ThreadPoolExecutor() as executor:
            responses = list(executor.map(fetch_response, prompts))

        def remove_emphasis(output):
            return output.replace("**", "")
            # Extract the individual responses

        other_summary = remove_emphasis(responses[0]['choices'][0]['message']['content'].strip())

        # Call to facility_threat_profile
        threatSummary, insightSummary, strategySummary = facility_threat_profile(sl, facility, facility_type, country,
                                                                                 Industry,
                                                                                 other_summary,
                                                                                 investment_statement, has_ir_plan_str,
                                                                                 ir_plan_never_tested_str,
                                                                                 ir_plan_date_str, connector_risk,
                                                                                 connector_status, connector_score)

        return JsonResponse({

            'other_summary': other_summary,

            'threatSummary': threatSummary,
            'insightSummary': insightSummary,
            'strategySummary': strategySummary
        })


@login_required()
def pha_report(request, cyberpha_id):
    organization_id_from_session = request.session.get('user_organization')

    users_in_organization = User.objects.filter(userprofile__organization__id=organization_id_from_session)

    cyberPHAHeader = tblCyberPHAHeader.objects.get(ID=cyberpha_id)
    cyberScenarios = tblCyberPHAScenario.objects.filter(CyberPHA=cyberpha_id, Deleted=0)
    average_impact_safety = cyberScenarios.aggregate(Avg('impactSafety'))['impactSafety__avg']
    average_impact_danger = cyberScenarios.aggregate(Avg('impactDanger'))['impactDanger__avg']
    average_impact_environment = cyberScenarios.aggregate(Avg('impactEnvironment'))['impactEnvironment__avg']
    average_impact_production = cyberScenarios.aggregate(Avg('impactProduction'))['impactProduction__avg']
    average_impact_finance = cyberScenarios.aggregate(Avg('impactFinance'))['impactFinance__avg']
    average_impact_data = cyberScenarios.aggregate(Avg('impactData'))['impactData__avg']
    average_impact_reputation = cyberScenarios.aggregate(Avg('impactReputation'))['impactReputation__avg']
    average_impact_regulation = cyberScenarios.aggregate(Avg('impactRegulation'))['impactRegulation__avg']
    average_impact_supply = cyberScenarios.aggregate(Avg('impactSupply'))['impactSupply__avg']
    total_cost_impact = cyberScenarios.aggregate(Sum('sle'))['sle__sum']
    total_cost_impact_low = cyberScenarios.aggregate(Sum('sle_low'))['sle_low__sum']
    total_cost_impact_high = cyberScenarios.aggregate(Sum('sle_high'))['sle_high__sum']

    control_effectiveness = math.ceil(get_overall_control_effectiveness_score(cyberpha_id))
    # Retrieve the top 5 most effective controls
    top_5_controls = MitreControlAssessment.objects.filter(cyberPHA=cyberpha_id).order_by('-effectiveness_percentage')[
                     :5]

    # Retrieve the bottom 5 least effective controls
    bottom_5_controls = MitreControlAssessment.objects.filter(cyberPHA=cyberpha_id).order_by(
        'effectiveness_percentage')[:5]

    overall_probability = math.ceil(calculate_overall_probability(cyberpha_id))

    return JsonResponse({'cyberPHAHeader': model_to_dict(cyberPHAHeader),
                         'scenarios': list(cyberScenarios.values()),
                         'average_impact_safety': average_impact_safety,
                         'average_impact_danger': average_impact_danger,
                         'average_impact_environment': average_impact_environment,
                         'average_impact_production': average_impact_production,
                         'average_impact_finance': average_impact_finance,
                         'average_impact_data': average_impact_data,
                         'average_impact_reputation': average_impact_reputation,
                         'average_impact_regulation': average_impact_regulation,
                         'average_impact_supply': average_impact_supply,
                         'control_effectiveness': control_effectiveness,
                         'top_5_controls': list(
                             top_5_controls.values('control__name', 'effectiveness_percentage', 'weighting')),
                         'bottom_5_controls': list(
                             bottom_5_controls.values('control__name', 'effectiveness_percentage', 'weighting')),
                         'overall_probability': overall_probability,
                         'total_cost_impact': total_cost_impact,
                         'total_cost_impact_low': total_cost_impact_low,
                         'total_cost_impact_high': total_cost_impact_high
                         })


def get_overall_control_effectiveness_score(cyberPHA_ID):
    """
    Calculate the Overall Control Effectiveness Score for a given cyberPHA_ID.

    Args:
        cyberPHA_ID (int): The ID of the cyberPHA security assessment.

    Returns:
        float: The Overall Control Effectiveness Score.
    """

    # Filter the assessments based on the given cyberPHA_ID
    assessments = MitreControlAssessment.objects.filter(cyberPHA=cyberPHA_ID)

    # Calculate the weighted sum of effectiveness percentages
    weighted_sum = assessments.aggregate(
        total=Sum(F('effectiveness_percentage') * F('weighting'))
    )['total'] or 0

    # Calculate the total weight
    total_weight = assessments.aggregate(total=Sum('weighting'))['total'] or 1

    # Calculate the Overall Control Effectiveness Score
    overall_score = weighted_sum / total_weight

    return overall_score


def calculate_overall_probability(cyberpha_id):
    # Retrieve all scenarios related to the given assessment
    cyberScenarios = tblCyberPHAScenario.objects.filter(CyberPHA=cyberpha_id, Deleted=0)

    # Start with the probability that none of the scenarios occur
    probability_none_occur = 1.0

    # Iterate over each scenario and update the probability
    for scenario in cyberScenarios:
        probability_value = int(scenario.probability.strip('%'))
        probability_none_occur *= (1 - probability_value / 100.0)

    # Calculate the overall probability that at least one scenario occurs
    overall_probability = 1 - probability_none_occur

    return overall_probability * 100  # Return as a percentage


@login_required()
def phascenarioreport(request):
    cyberPHAID = request.POST.get('cyberPHAID')
    cyberPHAHeader = tblCyberPHAHeader.objects.get(ID=cyberPHAID)
    cyberScenarios = tblCyberPHAScenario.objects.filter(CyberPHA=cyberPHAID, Deleted=0)

    average_impact_safety = cyberScenarios.aggregate(Avg('impactSafety'))['impactSafety__avg']
    average_impact_danger = cyberScenarios.aggregate(Avg('impactDanger'))['impactDanger__avg']
    average_impact_environment = cyberScenarios.aggregate(Avg('impactEnvironment'))['impactEnvironment__avg']
    average_impact_production = cyberScenarios.aggregate(Avg('impactProduction'))['impactProduction__avg']
    average_impact_finance = cyberScenarios.aggregate(Avg('impactFinance'))['impactFinance__avg']
    average_impact_data = cyberScenarios.aggregate(Avg('impactData'))['impactData__avg']
    average_impact_reputation = cyberScenarios.aggregate(Avg('impactReputation'))['impactReputation__avg']
    average_impact_regulation = cyberScenarios.aggregate(Avg('impactRegulation'))['impactRegulation__avg']

    return render(request, 'phascenarioreport.html', {
        'cyberPHAHeader': cyberPHAHeader,
        'cyberScenarios': cyberScenarios,
        'average_impact_safety': average_impact_safety,
        'average_impact_danger': average_impact_danger,
        'average_impact_environment': average_impact_environment,
        'average_impact_production': average_impact_production,
        'average_impact_finance': average_impact_finance,
        'average_impact_data': average_impact_data,
        'average_impact_reputation': average_impact_reputation,
        'average_impact_regulation': average_impact_regulation
    })


@login_required()
def getSingleScenario(request):
    # Get the ID from the GET parameters
    scenario_id = request.GET.get('id')
    current_user = request.user

    # Try to retrieve the scenario with the given ID
    try:
        scenario = tblCyberPHAScenario.objects.get(ID=scenario_id)
    except ObjectDoesNotExist:
        return JsonResponse({'error': 'Scenario not found'}, status=404)
        # Check if the current user is the owner, moderator, read-only, or has full access

    current_user_profile = UserProfile.objects.get(user=current_user)
    scenario_user_profile = UserProfile.objects.get(user=scenario.userID)

    if scenario.userID == current_user:
        user_role = 'Scenario Owner'
    else:
        try:
            user_profile = UserProfile.objects.get(user=current_user)
            if user_profile.role_moderator:
                user_role = 'Scenario Moderator'
            elif user_profile.role_readonly:
                user_role = 'Read Only'
            elif current_user_profile.organization_id == scenario_user_profile.organization_id:
                user_role = 'Scenario Editor'
            else:
                return JsonResponse({'error': 'Access denied'}, status=403)
        except UserProfile.DoesNotExist:
            return JsonResponse({'error': 'User profile not found'}, status=404)

    # Convert the scenario to a dictionary
    scenario_dict = {
        'ID': scenario.ID,
        'CyberPHA': scenario.CyberPHA.ID,
        'Scenario': scenario.Scenario,
        'ThreatClass': scenario.ThreatClass,
        'ThreatAgent': scenario.ThreatAgent,
        'ThreatAction': scenario.ThreatAction,
        'Countermeasures': scenario.Countermeasures,
        'RiskCategory': scenario.RiskCategory,
        'Consequence': scenario.Consequence,
        'impactSafety': scenario.impactSafety,
        'impactDanger': scenario.impactDanger,
        'impactProduction': scenario.impactProduction,
        'impactFinance': scenario.impactFinance,
        'impactReputation': scenario.impactReputation,
        'impactEnvironment': scenario.impactEnvironment,
        'impactRegulation': scenario.impactRegulation,
        'impactData': scenario.impactData,
        'impactSupply': scenario.impactSupply,
        'justifySafety': scenario.justifySafety,
        'justifyLife': scenario.justifyLife,
        'justifyProduction': scenario.justifyProduction,
        'justifyFinancial': scenario.justifyFinancial,
        'justifyReputation': scenario.justifyReputation,
        'justifyEnvironment': scenario.justifyEnvironment,
        'env_contaminant': scenario.env_contaminant,
        'env_ecosystem': scenario.env_ecosystem,
        'env_contamination': scenario.env_contamination,
        'env_population': scenario.env_population,
        'env_wildlife': scenario.env_wildlife,
        'justifyRegulation': scenario.justifyRegulation,
        'justifyData': scenario.justifyData,
        'UEL': scenario.UEL,
        'uel_threat': scenario.uel_threat,
        'uel_vuln': scenario.uel_vuln,
        'uel_exposure': scenario.uel_exposure,
        'RRU': scenario.RRU,
        'SM': scenario.SM,
        'MEL': scenario.MEL,
        'RRM': scenario.RRM,
        'SA': scenario.SA,
        'MELA': scenario.MELA,
        'RRa': scenario.RRa,
        'recommendations': scenario.recommendations,
        'Deleted': scenario.Deleted,
        'timestamp': scenario.timestamp,
        'aro': scenario.aro,
        'sle': scenario.sle,
        'sle_low': scenario.sle_low,
        'sle_high': scenario.sle_high,
        'ale_median': scenario.ale_median,
        'ale_low': scenario.ale_low,
        'ale_high': scenario.ale_high,
        'ale': scenario.ale,
        'countermeasureCosts': scenario.countermeasureCosts,
        'control_recommendations': scenario.control_recommendations,
        'probability': scenario.probability,
        'risk_register': scenario.risk_register,
        'safety_hazard': scenario.safety_hazard,
        'sis_compromise': scenario.sis_compromise,
        'sis_outage': scenario.sis_outage,
        'risk_priority': scenario.risk_priority,
        'risk_status': scenario.risk_status,
        'risk_owner': scenario.risk_owner,
        'risk_response': scenario.risk_response,
        'risk_open_date': scenario.risk_open_date,
        'risk_close_date': scenario.risk_close_date,
        'control_effectiveness': scenario.control_effectiveness,
        'likelihood': scenario.likelihood,
        'frequency': scenario.frequency,
        'sl_a': scenario.sl_a,
        'dangerScope': scenario.dangerScope,
        'prodOutage': scenario.outage,
        'prodOutageDuration': scenario.outageDuration,
        'prodOutageCost': scenario.outageCost,
        'ai_bia_score': scenario.ai_bia_score,
        'exposed_system': scenario.exposed_system,
        'weak_credentials': scenario.weak_credentials,
        'compliance_map': scenario.compliance_map,
        'attack_tree_text': scenario.attack_tree_text,
        'scenario_status': scenario.scenario_status,
        'cost_projection': scenario.cost_projection,
        'user_role': user_role,
        'risk_rationale': scenario.risk_rationale,
        'risk_recommendation': scenario.risk_recommendation,
        'cost_justification': scenario.cost_justification,
        'asset_name': scenario.asset_name,
        'asset_purpose': scenario.asset_purpose,
        'asset_critical': scenario.asset_critical,
        'detection_time': scenario.detection_time,
        'incident_complexity': scenario.incident_complexity,
        'response_time': scenario.response_time,
        'malware': {
            'id': scenario.malware.id,
            'name': scenario.malware.name,
            'description': scenario.malware.description,
            'source_link': scenario.malware.source_link
        } if scenario.malware else None  # Serialize the Malware object
    }
    # Retrieve the related PHAControlList records
    control_list = []
    for control in scenario.controls.all():
        control_dict = {
            'ID': control.ID,
            'control': control.control,
            'reference': control.reference,
            'score': control.score
        }
        control_list.append(control_dict)

    # Add the control list to the scenario dictionary
    scenario_dict['controls'] = control_list
    has_controls = 1 if len(control_list) > 0 else 0

    # Add has_controls to the scenario_dict
    scenario_dict['has_controls'] = has_controls
    # Retrieve related consequences
    consequences = ScenarioConsequences.objects.filter(scenario=scenario)
    consequences_list = [{'consequence_text': consequence.consequence_text, 'is_validated': consequence.is_validated}
                         for consequence in consequences]
    scenario_dict['Consequences'] = consequences_list
    # Retrieve the related PHA_Safeguard records
    safeguards = PHA_Safeguard.objects.filter(scenario=scenario)
    safeguards_list = [
        {
            'id': safeguard.id,
            'safeguard_description': safeguard.safeguard_description,
            'safeguard_type': safeguard.safeguard_type
        } for safeguard in safeguards
    ]

    # Add the safeguards list to the scenario dictionary
    scenario_dict['safeguards'] = safeguards_list

    observations = PHA_Observations.objects.filter(scenario=scenario)
    observations_list = [
        {
            'id': observation.id,
            'observation_description': observation.observation_description,

        } for observation in observations
    ]

    # Add the safeguards list to the scenario dictionary
    scenario_dict['observations'] = observations_list

    scenario_instance = tblCyberPHAScenario.objects.get(ID=scenario_id)
    write_to_audit(
        user_id=request.user,
        user_action=f'Viewed Scenario',
        user_ip=get_client_ip(request),
        cyberPHAScenario=scenario_instance,
    )

    # Return the scenario as a JSON response
    return JsonResponse(scenario_dict)


# function to calculate overall control effectiveness for a given cyberpha
def calculate_effectiveness(cyberPHA_value):
    # Filter assessments by the given cyberPHA value
    assessments = MitreControlAssessment.objects.filter(cyberPHA=cyberPHA_value)

    # If there are no matching records, return 0
    if not assessments.exists():
        return 0

    # Calculate the numerator (sum of effectiveness_percentage multiplied by weighting)
    numerator = assessments.aggregate(
        total_effectiveness=Sum(F('effectiveness_percentage') * F('weighting'))
    )['total_effectiveness']

    # Calculate the denominator (sum of weighting)
    denominator = assessments.aggregate(total_weighting=Sum('weighting'))['total_weighting']

    # Calculate the weighted average
    if denominator:
        overall_effectiveness = numerator / denominator
    else:
        overall_effectiveness = 0

    return overall_effectiveness


def get_response(user_message):
    openai_model = get_api_key("OpenAI_Model")
    message = [
        {
            "role": "system",
            "content": "You are an expert and experienced process and safety engineer conducting a cybersecurity risk analysis for a cyberPHA (where P=Process, H=Hazards, A=Analysis) scenario related to industrial automation and control systems."
        },
        user_message
    ]

    response = openai.ChatCompletion.create(
        model=openai_model,
        messages=message,
        temperature=0.1,
        max_tokens=4000
    )

    return response['choices'][0]['message']['content']


def compliance_map_data(common_content):
    pinecone_query = "Compliance, regulations, cybersecurity, law, regulatory compliance"
    retrieved_chunks = query_index(pinecone_query)
    summarized_chunks = get_summarized_chunks(retrieved_chunks)
    documents_context = "\n\n".join(summarized_chunks)

    user_message = {
        "role": "user",
        "content": f""""
            {common_content}. 
            
            Based on the provided information, please map the current OT security posture to a maximum of 10 of the MOST RELEVANT AND IMPORTANT industry regulatory compliance regulations for this organization in the given country. 
            When naming these regulations, use their official titles as recognized by the issuing bodies or as commonly used in official publications. 
            
            Extra guidance:
            - Be relevant to the given industry and type of facility
            - Be relevant to OT cybersecurity
            - Be relevant to country
            
            Formatting instruction for output: 
            
            Separate each item with '||' and use '>' to separate parts within an item. 
            Ensure each item's format is concise and can be easily parsed for display in an HTML table. 
            
           The precise format for each item must be as follows : 
            
            Concise Description of reason for the Compliance Concern, maximum 30 words > Official Compliance Reference > Internet URL ||
            
            Output only the line items with NO additional text, header, intro, or narrative. 
            Strive for consistency in the naming of compliance regulations to facilitate accurate parsing and display.
            Pinecone index data for extra context {documents_context}. 
            """
    }

    try:
        response = get_response(user_message)

        return response
    except Exception as e:
        return f"Error: {str(e)}"


def generate_recommendation_prompt(likelihood, adjustedRR, costs, probability, frequency, biaScore, scenario,
                                   cyberphaID):
    # Attempt to fetch the related CyberPHARiskTolerance record

    pinecone_query = "Cybersecurity Risk assessment, residual risk, threats, cybersecurity, hazops, cyberpha, risk mitigation, risk scoring"
    retrieved_chunks = query_index(pinecone_query)
    summarized_chunks = get_summarized_chunks(retrieved_chunks)
    documents_context = "\n\n".join(summarized_chunks)

    prompt = f"""
        Pinecone indexed data for reference: {documents_context}
        Given the cybersecurity risk assessment results for a given OT cybersecurity scenario described as: {scenario}.
        Likelihood of occurrence: {likelihood}%
        Adjusted residual risk: {adjustedRR}
        Estimated costs (low|medium|high): {costs}
        Probability of a targeted attack being successful: {probability}%
        Annual loss event frequency (as defined by FAIR): {frequency}
        Business impact score: {biaScore}

        Provide a response structured exactly as follows:

        1. Recommendation: [Insert recommendation here based on the assessment results. For example, Mitigate Risk, Accept Risk, Manage Risk, or another recommendation depending on results]
        2. Rationale: [Insert a concise, executive-level rationale here. Consider the likelihood of occurrence (values under 10 should be considered negligible), adjusted residual risk, business impact score, estimated costs, annual threat event frequency, probability of a targeted attack being successful, and risk tolerance levels. Structure your rationale in clear, bullet-pointed explanations.]
        """
    return prompt


def get_openai_recommendation(prompt):
    openai.api_key = get_api_key("openai")  # Ensure this retrieves your OpenAI API key correctly
    model = get_api_key("OpenAI_Model")

    messages = [
        {"role": "system",
         "content": "You are an expert in OT cybersecurity risk assessment. Provide a recommendation based on the analysis."},
        {"role": "user", "content": prompt}
    ]

    response = openai.ChatCompletion.create(
        model=model,
        messages=messages,
        temperature=0.2,
        max_tokens=600
    )

    if response.choices and len(response.choices) > 0:
        # Assuming the response follows the structured format: "1. Recommendation: ... 2. Rationale: ..."
        text = response.choices[0].message['content'].strip()
        text = text.replace("**", "")
        # Splitting the response into Recommendation and Rationale parts
        parts = text.split("2. Rationale:")
        recommendation = parts[0].replace("1. Recommendation:", "").strip()
        rationale = parts[1].strip() if len(parts) > 1 else ""

        # Structuring the response for JSON
        structured_response = {
            "Recommendation": recommendation,
            "Rationale": rationale
        }

    else:
        structured_response = {"Recommendation": "No recommendation could be generated.", "Rationale": ""}

    return structured_response


def exalens_get_device_incident(ipaddress, exalens_api_key, exalens_ip_address, exalens_client_id):
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    incident_url = f"https://{exalens_ip_address}/api/thirdparty/incident/target_ip/{ipaddress}?incident=1"
    headers = {
        'x-client-id': exalens_client_id,
        'x-api-key': exalens_api_key
    }
    try:
        response = requests.get(incident_url, headers=headers, verify=False)

        if response.status_code == 200:
            # Directly parse the response as JSON, which is expected to be a list of dictionaries
            incident_data = response.json()
            if not incident_data:  # Check if the list is empty
                return 'No incidents found for the given IP address.'

            # Filter the incident data to keep only the required fields
            filtered_incidents = [
                {
                    "detection_summary": incident.get("detection_summary"),
                    "risk_score": incident.get("risk_score"),
                    "kill_chain": incident.get("detection_artifacts", {}).get("kill_chain"),
                    "mitre_attack": incident.get("detection_artifacts", {}).get("mitre_attack")
                }
                for incident in incident_data
            ]

            # Convert the filtered incident data to a single string
            incidents_str = json.dumps(filtered_incidents, indent=2)
            return incidents_str
        else:
            # Handle non-200 responses
            return f'Failed to fetch incidents, status code: {response.status_code}'
    except RequestException as e:
        # Handle exceptions from the requests library
        return f'Error occurred: {str(e)}'


def format_currency(value):
    if value >= 1_000_000:
        return f"${value // 1_000_000}m"
    elif value >= 1_000:
        return f"${value // 1_000}k"
    else:
        return f"${value}"


def scenario_analysis_report(request):
    # Get parameters from the request
    scenario = request.GET.get('scenario')
    likelihood = int(request.GET.get('likelihood'))
    adjustedRR = request.GET.get('adjustedRR')
    costs = request.GET.get('costs')

    probability = request.GET.get('probability')
    frequency = request.GET.get('frequency')
    biaScore = int(request.GET.get('biaScore'))
    projection = request.GET.get('projection')
    cost_projection_justification = request.GET.get('cost_projection_justification')
    control_effectiveness = request.GET.get('control_effectiveness')
    recommendations = request.GET.get('recommendations')
    select_level = request.GET.get('select_level')
    scenario_compliance_data = request.GET.get('scenario_compliance_data')
    rationale = json.loads(request.GET.get('rationale'))
    CyberPHAID = request.GET.get('CyberPHAID')
    asset = request.GET.get('asset')
    asset_purpose = request.GET.get('asset_purpose')
    # Retrieve the CyberPHAHeader and related Facility data
    cyber_pha_header = get_object_or_404(tblCyberPHAHeader, pk=CyberPHAID)
    facility = get_object_or_404(Facility, pk=cyber_pha_header.facility_id)
    facility_type = get_object_or_404(FacilityType, pk=facility.type_id)

    # Path to the logo image
    logo_path = os.path.join('static/images', '65C8D0 - Light Blue-2.png')

    # Create a PDF buffer
    buffer = BytesIO()

    class CustomDocTemplate(BaseDocTemplate):
        def __init__(self, *args, **kwargs):
            super().__init__(*args, **kwargs)
            self.page_count = 0
            frame = Frame(self.leftMargin, self.bottomMargin, self.width, self.height, id='normal')
            template = PageTemplate(id='test', frames=frame, onPage=self.add_page_number_and_footer)
            self.addPageTemplates(template)

        def add_page_number_and_footer(self, canvas, doc):
            # Add page number
            page_num_text = f"Page {doc.page}"
            canvas.drawRightString(200 * mm, 10 * mm, page_num_text)
            # Add "Confidential" at the center bottom
            canvas.drawCentredString(self.leftMargin + self.width / 2, 10 * mm, "Confidential")

    # Create a document template
    doc = CustomDocTemplate(buffer, pagesize=letter)
    elements = []

    # Styles for paragraphs
    styles = getSampleStyleSheet()
    title_style = styles['Title']
    heading_style = styles['Heading1']

    helvetica_10 = ParagraphStyle(
        'Helvetica10',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
    )

    # Title page
    elements.append(Spacer(1, 2 * inch))
    elements.append(Paragraph("CyberPHA Scenario Report", title_style))
    elements.append(Spacer(1, 2 * inch))
    elements.append(Image(logo_path, width=180, height=150))
    elements.append(PageBreak())

    # Executive summary page
    elements.append(Paragraph("Executive Summary", heading_style))
    elements.append(Spacer(1, 0.4 * inch))

    summary_text = (f"The following scenario has been analyzed for CyberPHA: {cyber_pha_header.title} "
                    f"for {facility.name}, a {facility_type.FacilityType} at {facility.address}.")

    elements.append(Paragraph(summary_text, helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))
    elements.append(Paragraph("The scenario under analysis is:", helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Add the scenario text inside a box
    scenario_paragraphs = scenario.split('\n\n')
    scenario_content = [Paragraph(paragraph, helvetica_10) for paragraph in scenario_paragraphs]

    table_data = [[Paragraph(scenario_paragraphs[0], helvetica_10)]]
    for para in scenario_paragraphs[1:]:
        table_data.append([Spacer(1, 0.2 * inch)])
        table_data.append([Paragraph(para, helvetica_10)])

    table = Table([[table_data]], colWidths=[doc.width])
    table.setStyle(TableStyle([
        ('BOX', (0, 0), (-1, -1), 0.5, colors.black),
        ('LEFTPADDING', (0, 0), (-1, -1), 12),
        ('RIGHTPADDING', (0, 0), (-1, -1), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 12),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
    ]))

    elements.append(table)

    elements.append(PageBreak())

    # Scenario Event Costs page
    elements.append(Paragraph("Scenario Event Costs", heading_style))
    elements.append(Spacer(1, 0.4 * inch))
    elements.append(
        Paragraph("The financial impacts of the scenario have been estimated as follows below:", helvetica_10))
    elements.append(Spacer(1, 0.4 * inch))

    # Extract costs values
    cost_values = [val.strip() for val in costs.split('|')]
    cost_table_data = [
        ['Best Case', 'Most Likely Case', 'Worst Case'],
        [cost_values[0], cost_values[1], cost_values[2]]
    ]
    cost_table = Table(cost_table_data, hAlign='CENTER')
    cost_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('FONTNAME', (0, 0), (-1, -1), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (1, 0), (-1, -1), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('TOPPADDING', (0, 0), (-1, 0), 12),
        ('BOX', (0, 0), (-1, -1), 0.5, colors.black),
        ('INNERGRID', (0, 0), (-1, -1), 0.5, colors.black),
    ]))
    elements.append(cost_table)
    elements.append(Spacer(1, 0.2 * inch))

    # Generate the line graph for projection values using ReportLab
    projection_values = [int(val) for val in projection.split('|')]
    months = [f"Month {i + 1}" for i in range(12)]

    drawing = Drawing(400, 250)
    line_chart = HorizontalLineChart()
    line_chart.x = 50
    line_chart.y = 50
    line_chart.height = 125
    line_chart.width = 300
    line_chart.data = [projection_values]
    line_chart.categoryAxis.categoryNames = months
    line_chart.valueAxis.valueMin = 0
    line_chart.valueAxis.valueMax = max(projection_values) * 1.1
    line_chart.valueAxis.valueStep = max(projection_values) // 10

    line_chart.valueAxis.labelTextFormat = format_currency
    line_chart.categoryAxis.labels.angle = 45
    line_chart.categoryAxis.labels.boxAnchor = 'n'
    for y in range(int(line_chart.valueAxis.valueMin), int(line_chart.valueAxis.valueMax) + 1,
                   int(line_chart.valueAxis.valueStep)):
        drawing.add(Line(line_chart.x, line_chart.y + (line_chart.height * (y - line_chart.valueAxis.valueMin) / (
                line_chart.valueAxis.valueMax - line_chart.valueAxis.valueMin)),
                         line_chart.x + line_chart.width,
                         line_chart.y + (line_chart.height * (y - line_chart.valueAxis.valueMin) / (
                                 line_chart.valueAxis.valueMax - line_chart.valueAxis.valueMin)),
                         strokeColor=colors.lightgrey))

    drawing.add(line_chart)

    title = String(200, 180, "Monthly Incident Cost Projection", fontName='Helvetica-Bold', fontSize=14,
                   textAnchor='middle')
    drawing.add(title)

    elements.append(drawing)
    elements.append(Spacer(1, 0.4 * inch))

    # Cost projection justification
    elements.append(Paragraph("Cost Projection Justification:", helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Split the cost_projection_justification into lines and remove empty lines
    justification_bullets = [line.strip() for line in cost_projection_justification.split('\n') if line.strip()]
    for bullet in justification_bullets:
        elements.append(Paragraph(f"• {bullet}", helvetica_10))
        elements.append(Spacer(1, 0.1 * inch))

    elements.append(PageBreak())
    elements.append(Paragraph("Risk Outcomes", heading_style))
    elements.append(Spacer(1, 0.4 * inch))

    # Likelihood text
    likelihood = int(likelihood)
    likelihood_text = "The likelihood of this scenario occurring is considered to be: "
    if likelihood < 25:
        likelihood_text += "<b>Low</b>"
    elif 25 <= likelihood <= 40:
        likelihood_text += "<b>Low/Medium</b>"
    elif 41 <= likelihood <= 65:
        likelihood_text += "<b>Medium</b>"
    elif 66 <= likelihood <= 80:
        likelihood_text += "<b>Medium/High</b>"
    elif 81 <= likelihood <= 95:
        likelihood_text += "<b>High</b>"
    elif 96 <= likelihood <= 100:
        likelihood_text += "<b>Very High</b>"

    elements.append(Paragraph(likelihood_text, helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Residual Risk text
    residual_risk_text = f"Residual Risk: <b>{adjustedRR}</b>"
    elements.append(Paragraph(residual_risk_text, helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Probability text
    probability_text = f"The probability of an attack based on this scenario being successful is estimated at: <b>{probability}</b>"
    elements.append(Paragraph(probability_text, helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Overall Business Impact text
    biaScore = int(biaScore)
    bia_text = "Overall Business Impact: "
    if biaScore < 25:
        bia_text += "<b>Low</b>"
    elif 25 <= biaScore <= 40:
        bia_text += "<b>Low/Medium</b>"
    elif 41 <= biaScore <= 65:
        bia_text += "<b>Medium</b>"
    elif 66 <= biaScore <= 80:
        bia_text += "<b>Medium/High</b>"
    elif 81 <= biaScore <= 95:
        bia_text += "<b>High</b>"
    elif 96 <= biaScore <= 100:
        bia_text += "<b>Very High</b>"

    elements.append(Paragraph(bia_text, helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Rationale text
    elements.append(Paragraph("Overall Recommendation: ", helvetica_10))
    elements.append(Paragraph(f"<b>{rationale['Recommendation']}</b>", helvetica_10))
    elements.append(Spacer(1, 0.2 * inch))

    # Bullet points for Rationale details
    rationale_points = rationale['Rationale'].split('\n')
    for point in rationale_points:
        elements.append(Paragraph(f"• {point}", helvetica_10))
        elements.append(Spacer(1, 0.1 * inch))

    # Scenario Recommendations page
    elements.append(PageBreak())
    elements.append(Paragraph("Scenario Recommendations", heading_style))
    elements.append(Spacer(1, 0.4 * inch))

    # Prepare the table data with word-wrapping
    recommendations_list = recommendations.split('\n')
    table_data = [['#', 'Recommendation', 'NIST Reference']]

    styles = getSampleStyleSheet()
    normal_style = styles['Normal']

    for rec in recommendations_list:
        if rec.strip():
            parts = rec.split(' [')
            number_and_recommendation = parts[0].split('. ', 1)
            nist_reference = parts[1].strip(']')
            table_data.append([
                number_and_recommendation[0],
                Paragraph(number_and_recommendation[1], normal_style),
                Paragraph(nist_reference, normal_style)
            ])

    # Define the table with specific column widths
    col_widths = [0.08 * doc.width, 0.70 * doc.width, 0.22 * doc.width]
    recommendations_table = Table(table_data, colWidths=col_widths)

    # Apply table style
    recommendations_table.setStyle(TableStyle([
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 12),
    ]))

    elements.append(recommendations_table)
    elements.append(Spacer(1, 0.4 * inch))

    elements.append(PageBreak())
    elements.append(Paragraph("Appendix: Regulatory Compliance Requirements", heading_style))
    elements.append(Spacer(1, 0.4 * inch))
    elements.append(Paragraph(
        "The following table displays the regulatory compliance requirements considered to be most relevant to the given scenario.",
        helvetica_10))
    elements.append(Spacer(1, 0.4 * inch))

    # Prepare the table data
    compliance_data_list = scenario_compliance_data.split(' ||\n')
    table_data = [['Issue', 'Framework', 'Link']]

    for data in compliance_data_list:
        if data.strip():
            parts = data.split(' > ')
            issue = parts[0]
            framework = parts[1]
            link = parts[2]
            table_data.append([
                Paragraph(issue, normal_style),
                Paragraph(framework, normal_style),
                Paragraph(f'<link href="{link}">{link}</link>', normal_style)
            ])

    # Define the table with specific column widths
    col_widths = [0.33 * doc.width, 0.33 * doc.width, 0.33 * doc.width]
    compliance_table = Table(table_data, colWidths=col_widths)

    # Apply table style
    compliance_table.setStyle(TableStyle([
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
        ('ALIGN', (0, 0), (0, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
        ('TOPPADDING', (0, 0), (-1, -1), 12),
    ]))

    elements.append(compliance_table)
    elements.append(Spacer(1, 0.4 * inch))

    # Build the PDF
    doc.build(elements)

    # Get the PDF content
    buffer.seek(0)
    pdf_content = buffer.getvalue()
    response = HttpResponse(pdf_content, content_type='application/pdf')
    response['Content-Disposition'] = 'attachment; filename="scenario_analysis_report.pdf"'
    return response


@login_required
def scenario_analysis(request):
    # Log the user activity
    user_profile = UserProfile.objects.get(user=request.user)

    # Check if the user has exceeded the maximum scenario analysis count
    if user_profile.current_scenario_count >= user_profile.max_scenario_count:
        return JsonResponse({
            'error': 'User has reached the maximum number of scenario assessments'
        }, status=403)  # 403 Forbidden or another appropriate status code

    if request.method == 'GET':
        industry = request.GET.get('industry')
        facility_type = request.GET.get('facility_type')
        scenario = request.GET.get('scenario')
        threatSource = request.GET.get('threatsource')
        safetyimpact = request.GET.get('safety')
        safety_hazard = request.GET.get('safety_hazard')
        lifeimpact = request.GET.get('life')
        productionimpact = request.GET.get('production')
        production_outage = request.GET.get('production_outage')
        production_outage_length = request.GET.get('production_outage_length')
        reputationimpact = request.GET.get('reputation')
        environmentimpact = request.GET.get('environment')
        regulatoryimpact = request.GET.get('regulatory')
        dataimpact = request.GET.get('data')
        supplyimpact = request.GET.get('supply')
        malware = request.GET.get('malware')

        country = request.GET.get('country')

        financial = request.GET.get('financial')
        cyberPHAID = request.GET.get('cpha')
        exposed_system = request.GET.get('exposed_system')
        weak_credentials = request.GET.get('weak_credentials')
        # Retrieve the string of validated consequences
        validated_consequences_str = request.GET.get('validated_consequences', '')
        physical_safeguards_str = request.GET.get('physical_safeguards', '')
        observations = request.GET.get('observations')
        force_resubmit = request.GET.get('force_resubmit', 'false').lower() == 'true'
        safetyJustification = request.GET.get('safetyJustification')
        supplychainJustification = request.GET.get('supplychainJustification')
        environmentJustification = request.GET.get('environmentJustification')
        dangerJustification = request.GET.get('dangerJustification')
        targetAsset = request.GET.get('targetAsset')
        targetAssetPurpose = request.GET.get('targetAssetPurpose')
        connectionFlag = request.GET.get('partner_connection')

        # Concatenate all GET parameters to form a string
        values_str = '|'.join([request.GET.get(param, '') for param in request.GET])
        # Generate SHA-256 hash
        hash_value = hashlib.sha256(values_str.encode()).hexdigest()

        # Attempt to find a matching record
        existing_record = UserScenarioHash.objects.filter(
            user=request.user,
            cyberphaID=request.GET.get('cpha', 0),
            hash_value=hash_value
        ).exists()

        if existing_record and not force_resubmit:
            # If a matching record is found, inform the user and exit
            return JsonResponse({'message': 'Scenario parameters have not changed.'}, status=200)

        if not existing_record:
            UserScenarioHash.objects.create(
                user=request.user,
                cyberphaID=request.GET.get('cpha', ''),
                hash_value=hash_value
            )
        # Split the string into a list, using semicolon as the separator
        validated_consequences_list = validated_consequences_str.split(';')

        cyber_pha_header = tblCyberPHAHeader.objects.get(ID=cyberPHAID)

        has_incident_response_plan = cyber_pha_header.has_incident_response_plan
        plan_last_tested_date = cyber_pha_header.plan_last_tested_date
        plan_never_tested = cyber_pha_header.plan_never_tested
        exalens_api_key = cyber_pha_header.exalens_api
        exalens_ip_address = cyber_pha_header.exalens_ip
        exalens_client_id = cyber_pha_header.exalens_client

        # exalens_incidents = exalens_get_device_incident(targetAsset, exalens_api_key, exalens_ip_address,
        #                                                exalens_client_id)

        # Convert plan_last_tested_date to string format if it's not None
        plan_last_tested_date_str = plan_last_tested_date.strftime('%Y-%m-%d') if plan_last_tested_date else None

        ####### Audit Write ########
        write_to_audit(
            request.user,
            f'Executed a scenario analysis for {cyber_pha_header.title} and scenario: {scenario}',
            get_client_ip(request)
        )
        ###########################

        # Get the assessment id from the tblCyberPHAHeader instance
        assessment_id = cyber_pha_header.assessment
        last_assessment_score = cyber_pha_header.last_assessment_score
        last_assessment_summary = cyber_pha_header.last_assessment_summary

        investments = CyberSecurityInvestment.objects.filter(cyber_pha_header=cyber_pha_header).values(
            'investment_type', 'vendor_name', 'product_name', 'cost', 'date'
        )

        if investments.exists():
            # Generate the text statement for investments
            investment_statement = "OT Cybersecurity Investments listed here:\n"
            for idx, investment in enumerate(investments, start=1):
                investment_date = investment['date'].strftime('%Y-%m-%d') if investment['date'] else 'N/A'
                investment_statement += f"{idx}: Investment Type: {investment['investment_type']}, Vendor: {investment['vendor_name']}, Product: {investment['product_name']}, Cost: {investment['cost']}, Investment date: {investment_date}.\n"
        else:
            # Return an empty string if there are no investment records
            investment_statement = ""

        employees_on_site = cyber_pha_header.EmployeesOnSite
        # get the compliance summary for the cyberpha
        compliance = cyber_pha_header.complianceSummary
        net_profit_margin = cyber_pha_header.npm
        cost_op_hour = cyber_pha_header.coho
        annual_revenue = cyber_pha_header.annual_revenue

        asset_data = {}
        if connectionFlag == 1:
            darktrace_host = cyber_pha_header.darktrace_client
            darktrace_public_token = cyber_pha_header.darktrace_public_api
            darktrace_private_token = cyber_pha_header.darktrace_private_api

            # Fetch asset details from Darktrace
            asset_data = darktrace_asset_detail(targetAsset, darktrace_host, darktrace_public_token,
                                                darktrace_private_token)

        if connectionFlag == 2:
            exalens_client = cyber_pha_header.exalens_client
            exalens_api = cyber_pha_header.exalens_api
            exalens_url = cyber_pha_header.exalens_ip

            # Fetch asset details from Darktrace
            asset_data = exalens_asset_detail(targetAsset, exalens_url, exalens_api,
                                              exalens_client)

        if connectionFlag == 0:
            asset_data = 'No detailed asset data available'

        if malware != '':
            malware_data_query = f"Details about {malware}"
            retrieved_chunks = query_index(malware_data_query)
            malware_data = get_summarized_chunks(retrieved_chunks)
        else:
            malware_data = ''

        openai.api_key = get_api_key('openai')
        # Define the common part of the user message
        pinecone_query = "FAIR, risk assessment, 62443, cybersecurity, hazops, cyberpha, incident, event costs, costs, cost of a security incident "
        retrieved_chunks = query_index(pinecone_query)
        summarized_chunks = get_summarized_chunks(retrieved_chunks)
        documents_context = "\n\n".join(summarized_chunks)

        common_content = f"""
        Pinecone index data for reference: {documents_context}.
        Act as both an Insurance Actuary and an OT Cybersecurity HAZOPS Risk Expert. You're tasked with analyzing a specific scenario for a facility in the industry sector, located in a particular country. Your analysis should cover various risk outcomes based on the detailed context provided.

        Scenario Details:

        Facility Type & Industry: A {facility_type} in the {industry} industry, located in {country}.
        Scenario Overview: {scenario}.
        Critical System Exposures:
        Internet Exposure: Systems with public IP addresses: {exposed_system}.
        Affected asset: {targetAsset},
        Purpose of Affected Asset: {targetAssetPurpose}
        Asset_Data : {asset_data},
        Malware data: {malware_data},
        Credential Security: Systems with weak/default credentials: {weak_credentials}.
        OT Cybersecurity Incident Response:
        Presence of an Incident Response Plan: {has_incident_response_plan}.
        Last Tested Date: {plan_last_tested_date_str}.
        Estimated Business Impact Scores:
        Safety: {safetyimpact}, Life Danger: {lifeimpact}
        Production: {productionimpact}, Production Outage: {production_outage} (Duration: {production_outage_length} hours)
        Reputation: {reputationimpact}, Environmental: {environmentimpact}
        Regulatory Compliance: {regulatoryimpact}, Supply Chain: {supplyimpact}
        Data & Intellectual Property: {dataimpact}
        Estimated Current OT Cybersecurity Controls:
         OT Cybersecurity Control Effectiveness score: {last_assessment_score}/100.  OT Cybersecurity Control Effectiveness Summary: {last_assessment_summary}.
        Physical Security:
        Physical Safeguards: {physical_safeguards_str} (Assumed effective)
        Vulnerability Observations: {observations}
        Investment Statement: {investment_statement}
       
        """

        # Define the refined user messages
        user_messages = [

            {
                "role": "user",
                "content": f" {common_content} .Task: Given the scenario described and publicly reported cybersecurity incidents over the time scale from 5 years ago to the current day , estimate the likelihood (as a whole number percentage) of the given scenario occurring against a {facility_type} in {country}. Answer with a whole number. Do NOT include any other words, sentences, or explanations."
            },
            {
                "role": "user",
                "content": f"{common_content}. The FAIR (Factor Analysis of Information Risk) methodology, which is used to quantify and manage risk, defines residual risk as the amount of risk that remains after controls are applied. Task: Given the scenario described and publicly reported cybersecurity incidents over the time scale from 5 years ago to the current day, assess the cybersecurity residual risk after all recommended controls and physical security controls have been implemented and are assumed to be at least 75% effective. Give an estimated residual risk rating based on the FAIR methodology from the following options: Very Low, Low, Low/Medium, Medium, Medium/High, High, Very High.  Provide ONLY one of the given risk ratings without any additional text or explanations."
            },

            {
                "role": "user",
                "content": f"{common_content}.  Consequences of the scenario are assumed to be as follows: {validated_consequences_list}. TASK: Read and comply with all instructions that follow. Provide an estimate of the DIRECT COSTS of a single loss event (SLE) in US dollars. Provide three cost estimates: best case, worst case, and most likely case. Output these estimates as integers in the specific format: 'low|medium|high', where '|' is the delimiter. Ensure that your estimates are realistic, taking into account recent (no older than 3 years) industry reports about the cost of cybersecurity incidents, and tailored to the scenario, focusing solely on relevant Direct costs such as incident response, remediation, legal fees, notification costs, regulatory fines, compensations, and increased insurance premiums. The financial impact for this scenario is estimated as {financial}/10 in the business impact analysis. (IMPORTANT take into account the  OT Cybersecurity investments that have been made). IMPORTANT: Respond with only three positive integers, representing the low, medium, and high estimates, in the exact order and format specified: 'low|medium|high'. Do not include any additional text, explanations, or commentary."
            },

            {
                "role": "user",
                "content": f"""
                {common_content}
                TASK: Assess the likelihood of a successful targeted attack on the specified facility, given the described scenario and cybersecurity investments. Consider the effectiveness of implemented controls and any vulnerabilities due to the facility's internet exposure or weak/default credentials.
                Provide ONLY the estimated probability of such an attack being successful, expressed as a whole number percentage (e.g., 25%). Do NOT include any additional text, explanations, or commentary.
                """
            },
            # {
            #    "role": "user",
            #    "content": f"{common_content}. Provide an estimate of the annual Threat Event Frequency (TEF) as defined by the Open FAIR Body of Knowledge. TEF is the probable frequency, within a given timeframe, that a threat agent will act against an asset. It reflects the number of attempts by a threat actor to cause harm, regardless of whether these attempts are successful. Your response should reflect an integer or a decimal value representing the estimated number of times per year such a threat event is expected to occur. IMPORTANT: Respond only with the frequency value as an integer or a decimal, without including any additional words, sentences, or explanations. This value should strictly represent the estimated annual occurrence rate of the threat event as per the given scenario."
            #
            # },
            {
                "role": "user",
                "content": f"{common_content}. Provide an estimate of the annual Loss Event Frequency (LEF) as defined by the Open FAIR Body of Knowledge. LOSS Event Frequency is defined by the technical Open FAIR Body of Knowledge as: The probable frequency, within a given timeframe, that a threat agent will SUCCESSFULLY act against an asset. Your response should reflect an integer or a decimal value representing the estimated number of times per year such a threat event is expected to occur. IMPORTANT: Respond only with the frequency value as an integer or a decimal, without including any additional words, sentences, or explanations. This value should strictly represent the estimated annual occurrence rate of the threat event as per the given scenario."

            },
            {
                "role": "user",
                "content": f"{common_content}. TASK: Hypothesize the business impact score from 0 to 100 in the event of a successful attack resulting in the given scenario. Consequences of the scenario are given as follows: {validated_consequences_list}. A score of 1 would mean minimal business impact while a score of 100 would indicate catastrophic business impact without the ability to continue operations. Your answer should be given as an integer. Do NOT include any other words, sentences, or explanations."
            },

            {
                "role": "user",
                "content": f"Pinecone index data for reference: {documents_context}.I am a CISO at a {industry} company with approximately {employees_on_site} employees, operating primarily in {country}. We are assessing our cybersecurity posture and need to estimate the potential costs associated with a {scenario} that has consequences of {validated_consequences_list}."
                           "Your Task: Given the scenario, please provide an estimate of the direct and indirect costs we might incur, including but not limited to:"
                           "1. Immediate Response Costs: Costs associated with the initial response to the incident, such as emergency IT support, forensic analysis, and legal consultations."
                           "2. Remediation Costs: Expenses related to remediating the cybersecurity breach, including software updates, hardware replacements, and strengthening of security measures."
                           "3. Regulatory and Compliance Costs: Potential fines and penalties for non-compliance with relevant data protection and privacy regulations, as well as costs associated with compliance audits and reporting requirements post-incident."
                           "4. Reputation and Brand Impact: Estimated impact on our brand and customer trust, potentially leading to loss of business and decreased revenue."
                           "5. Operational Disruption: Costs associated with operational disruptions or downtime, including loss of productivity and impact on service delivery."
                           "6. Legal and Settlement Costs: Expenses related to legal actions taken against the company and any settlements or compensations paid out to affected parties."
                           "7. Long-term Costs: Any long-term costs such as increased insurance premiums, ongoing monitoring and security measures, and potential loss of intellectual property."
                           "INDUSTRY DATA AND ANALYTICS: Please consider the specifics of our industry, size, and the nature of the assets involved in this scenario to provide a comprehensive cost estimate. Please reference industry-specific data from the latest research and findings relating to the cost of a cybersecurity incident and data breach from Gartner, McKinsey, Dragos, Ponemon Institute, Verizon DBIR, Palo Alto Networks."
                           "OUTPUT INSTRUCTION: First, provide a 12-month direct cost projection for the scenario in the format: COST PROJECTION: Month1value|Month2value|...|Month12value. Each value must be an integer to represent a dollar value wth no other text or narrative included and should reflect a realistic, pragmatic monthly estimate. Then, provide a concise and conservative executive-level explanation summary, in 150 words or less, specific to the given scenario as JUSTIFICATION: <Your justification here>. Ensure the cost projection and justification are clearly separated by these keywords."
            },

            {
                "role": "user",
                "content": f"""{common_content}
                            You are an OT Cybersecurity expert. Consider the given scenario and all of the associated information.
                            Using the given details, generate a concise priority ordered and numbered bullet point list of OT/ICS cybersecurity action items that are readily implementable, as a check list, to PREVENT THE SCENARIO FROM OCCURRING ASSUMING IT HAS NOT YET OCCURRED. Each action item must be specific to the given scenario and represent a task that can be quickly and readily completed without major effort AND that would represent some risk mitigation. Each action item MUST be pragmatic and reasonable to accomplish within an OT network environment without extensive effort and could be implemented by an engineer rather than an IT security expert. For example Implement network segmentation is a major effort. Action items are intended to be akin to an immediate action drill. Identify which section of NIST 800-82 OR NIST CSF the action items most closely aligns with and include the relevant NIST reference in brackets at the end of each recommendation however action items MUST NOT SIMPLY BE QUOTES FROM THE STANDARDS. The output should strictly adhere to the following format:

                            Example Format:
                           1. Concise and pragmatic action item. [NIST Reference the action most closely relates to]
                           2. Another concise and pragmatic action item . [NIST Reference the action most closely relates to]

                            Following this example format, provide the actions in order of priority, specific to the given scenario without any additional narrative, description, advice, or guidance. The actions should be clear and easily parsable within an HTML page.
                            """

            },
            {
                "role": "user",
                "content": f"{common_content}. ISA-62443-3-2 describes five Security Levels SL-1, SL-2, SL-3, SL-4, SL-5. As an OT Cybersecurity risk analyst, assess which security level has been achieved and format the response exactly as follows without any additional text, narrative, characters or explanation: 'Security Level: [SL Value], Justification: [Concise Justification, no more than 20 words].'"
            }
        ]

        def get_response_safe(user_message):
            try:
                return get_response(user_message)
            except Exception as e:
                return f"Error: {str(e)}"

        # Initialize an empty list to store the responses
        responses = []

        # Use ThreadPoolExecutor to parallelize the first five API calls
        with concurrent.futures.ThreadPoolExecutor() as executor:
            # Submit all the tasks and get a list of futures
            futures = [executor.submit(get_response, msg) for msg in user_messages]
            # Collect the results in the order the futures were submitted
            responses = [future.result() for future in futures]

        # Now handle the compliance mapping separately
        compliance_data = compliance_map_data(common_content)
        responses.append(compliance_data)

        recommendation_prompt = generate_recommendation_prompt(
            likelihood=responses[0],
            adjustedRR=responses[1],
            costs=responses[2],
            probability=responses[3],
            frequency=responses[4],
            biaScore=responses[5],
            scenario=scenario,
            cyberphaID=cyberPHAID
        )

        # Get the recommendation from OpenAI
        rationale = get_openai_recommendation(recommendation_prompt)

        user_profile.current_scenario_count += 1
        user_profile.save()
        cost_projection_with_justification = responses[6]

        # Splitting the response to separate the cost projection and the justification
        cost_projection_parts = cost_projection_with_justification.split("COST PROJECTION: ")
        justification_parts = cost_projection_parts[1].split("JUSTIFICATION: ") if len(cost_projection_parts) > 1 else [
            "", "Justification not provided."]

        cost_projection = justification_parts[0].strip()
        cost_projection_justification = justification_parts[1].strip() if len(
            justification_parts) > 1 else "Justification not provided."

        ## pdf_buffer = scenario_analysis_report(
        ##    scenario=scenario,
        ##    likelihood=responses[0],
        ##    adjustedRR=responses[1],
        ##    costs=responses[2],
        ##    probability=responses[3],
        ##    frequency=responses[4],
        ##    biaScore=responses[5],
        ##    projection=cost_projection,
        ##    cost_projection_justification=cost_projection_justification,
        ##    control_effectiveness=last_assessment_score,
        ##    recommendations=responses[7],
        ##    select_level=responses[8],
        ##    scenario_compliance_data=responses[9],
        ##    rationale=rationale,
        ##    CyberPHAID=cyberPHAID,
        ##    asset=targetAsset,
        ##    asset_purpose=targetAssetPurpose,
        ##)
        ## pdf_base64 = base64.b64encode(pdf_buffer).decode('utf-8')
        # Return the responses as variables
        return JsonResponse({
            'likelihood': responses[0],
            'adjustedRR': responses[1],
            'costs': responses[2],
            'probability': responses[3],
            'frequency': responses[4],
            'biaScore': responses[5],
            'projection': cost_projection,
            'cost_projection_justification': cost_projection_justification,
            'control_effectiveness': last_assessment_score,
            'recommendations': responses[7],
            'select_level': responses[8],
            'scenario_compliance_data': responses[9],
            'rationale': rationale
            ## 'pdf': pdf_base64
        })


def prepare_controls_summary():
    controls = SecurityControls.objects.filter(framework='ISA 62443-2-1')
    summary = ""
    for control in controls:
        summary += f"Control: {control.Control}"
    return summary


def get_recommended_controls(scenario, threat_source):
    # Prepare the summary of controls
    controls_summary = prepare_controls_summary()

    # Construct the prompt for OpenAI
    prompt = f"Given these controls from the ISA 62443-2-1 standard:\n{controls_summary}\nWhat are the top 5 controls that would be most applicable for a scenario involving '{scenario}' with a threat source of '{threat_source}'? "

    # Make a call to OpenAI API
    try:
        response = openai.Completion.create(
            engine="text-davinci-002",  # Or the appropriate engine you are using
            prompt=prompt,
            max_tokens=250,  # Adjust as needed
            temperature=0.7  # Adjust for creativity level
        )
        return response.choices[0].text.strip()
    except Exception as e:
        # Handle exceptions (e.g., API errors)

        return None


def parse_recommendations(recommendations_str):
    # Regular expression to match the format
    pattern = re.compile(r'(\d+)\.\s(.*?)(\([^)]+\))')  # This regex captures any content within brackets

    matches = pattern.findall(recommendations_str)

    # Convert matches to a list of dictionaries
    parsed_data = [{
        'line_number': int(match[0]),
        'control_text': match[1].strip(),
        'reference': match[2]
    } for match in matches]

    return parsed_data


def scenario_vulnerability(request, scenario_id):
    try:
        scenario = tblCyberPHAScenario.objects.get(pk=scenario_id)
        vulnerabilities = vulnerability_analysis.objects.filter(scenario=scenario)
        asset_types = tblAssetType.objects.all()  # Fetch asset types

        if request.method == 'POST':
            form = VulnerabilityAnalysisForm(request.POST)
            if form.is_valid():
                form.save()
                # Redirect to the 'scenario_vulnerability' view with the scenario_id
                return redirect('OTRisk:scenario_vulnerability', scenario_id=scenario_id)
        else:
            form = VulnerabilityAnalysisForm()

        return render(request, 'OTRisk/vulnerability_table.html', {
            'scenario': scenario,
            'vulnerabilities': vulnerabilities,
            'form': form,
            'scenario_id_value': scenario_id,
            'asset_types': asset_types
        })
    except tblCyberPHAScenario.DoesNotExist:
        # Handle the scenario not found error, e.g., return a 404 page
        return render(request, 'no_vulnerabilities.html')


def add_vulnerability(request, scenario_id):
    scenario = get_object_or_404(tblCyberPHAScenario, pk=scenario_id)
    action = request.POST.get('action')
    if request.method == 'POST':
        if action == 'save':  # Check if Save button was clicked
            description = request.POST.get('description')
            asset_type_id = request.POST.get('asset_type')
            asset_name = request.POST.get('asset_name')
            cve = request.POST.get('cve')
            cve_detail = request.POST.get('cve_detail')

            # Create a new vulnerability instance and save it
            vulnerability = vulnerability_analysis(
                description=description,
                asset_name=asset_name,
                asset_type_id=asset_type_id,
                cve=cve,
                scenario=scenario,
                cve_detail=cve_detail
            )
            vulnerability.save()

        elif action == 'edit':  # Check if Update button was clicked
            vulnerability_id = request.POST.get('vulnerability_id')
            description = request.POST.get('description')
            asset_type_id = request.POST.get('asset_type')
            asset_name = request.POST.get('asset_name')
            cve = request.POST.get('cve')
            cve_detail = request.POST.get('cve_detail')

            # Get the existing vulnerability instance
            vulnerability = get_object_or_404(vulnerability_analysis, pk=vulnerability_id)

            # Update the fields and save
            vulnerability.description = description
            vulnerability.asset_type_id = asset_type_id
            vulnerability.asset_name = asset_name
            vulnerability.cve = cve
            vulnerability.cve_detail = cve_detail
            vulnerability.save()

    return redirect('OTRisk:scenario_vulnerability', scenario_id=scenario_id)


def get_asset_types(request):
    asset_types = tblAssetType.objects.all()
    asset_type_list = [{'id': asset_type.id, 'AssetType': asset_type.AssetType} for asset_type in asset_types]
    return JsonResponse({'asset_types': asset_type_list})


def generate_ppt(request):
    if request.method == "POST":

        # Extract data from POST parameters
        safety = request.POST.get('txtSafetySummary', '')
        chemical = request.POST.get('txtChemical', '')
        physical = request.POST.get('txtPhysical', '')
        other = request.POST.get('txtOther', '')
        compliance = request.POST.get('txtCompliance', '')
        facility = request.POST.get('FacilityName', '')
        threatSummary = request.POST.get('threatSummary', '')
        insightSummary = request.POST.get('insightSummary', '')
        strategySummary = request.POST.get('strategySummary', '')

        # Create a new PowerPoint presentation
        prs = Presentation()

        # For each section, add a slide and set its title and content
        sections = [
            (f"{facility} Safety Profile", safety),
            (f"{facility} Chemical Profile", chemical),
            (f"{facility} Physical Security Profile", physical),
            (f"{facility} OT Asset Profile", other),
            (f"{facility} Compliance Profile", compliance),
            (f"{facility} Threat Summary", threatSummary),
            (f"{facility} Security Insights", insightSummary),
            (f"{facility} Suggested Strategy", strategySummary)
        ]

        for title, content in sections:
            slide_layout = prs.slide_layouts[1]  # Using a title and content layout
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]  # Using the primary content placeholder

            title_shape.text = title
            title_shape.text_frame.paragraphs[0].font.name = 'Arial'
            title_shape.text_frame.paragraphs[0].font.size = Pt(18)

            content_shape.text = content
            for paragraph in content_shape.text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(10)

        # Save to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            prs.save(temp_file.name)
            temp_path = temp_file.name

        # Return a response to provide the file for download
        response = FileResponse(open(temp_path, 'rb'),
                                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename=report.pptx'

        # Remove the temporary file after sending it to the client
        os.remove(temp_path)
        filename = f"report_{uuid.uuid4()}.pptx"

        filepath = os.path.join(os.path.join(BASE_DIR, 'static'), filename)

        prs.save(filepath)
        download_url = urljoin(settings.STATIC_URL, filename)

        return JsonResponse({
            'status': 'success',
            'download_url': download_url
        })

    else:
        return JsonResponse({
            'status': 'error',
            'error': 'Invalid request method'
        })


def reformat_attack_tree(data):
    # Recursive function to adjust the format
    def adjust_node_format(node):
        new_node = {'name': node['Node']}
        if 'Children' in node and node['Children']:
            new_node['children'] = [adjust_node_format(child) for child in node['Children']]
        return new_node

    formatted_tree = {'name': data['Attack'], 'children': []}
    for node in data['Nodes']:
        formatted_tree['children'].append(adjust_node_format(node))

    return formatted_tree


@login_required
def analyze_scenario(request):
    openai_api_key = get_api_key('openai')
    openai.api_key = openai_api_key

    if request.method == 'POST':
        scenario = request.POST.get('scenario')
        attack_tree_drawn = request.POST.get('attackTreeDrawn') == 'true'
        attacker = request.POST.get('attacker')
        riskCategory = request.POST.get('riskCategory')
        attack_vector = request.POST.get('attack_vector')
        exposed_system = request.POST.get('exposed_system')
        weak_credentials = request.POST.get('weak_credentials')
        malware = request.POST.get('malware')

        # Fetch the organization_id from the user's profile
        try:
            user_profile = UserProfile.objects.get(user=request.user)
            organization_id = user_profile.organization.id
        except UserProfile.DoesNotExist:
            organization_id = None  # Or handle the lack of a profile as you see fit

        # Create a record in user_scenario_audit
        user_scenario_audit.objects.create(
            scenario_text=scenario,
            user=request.user,
            organization_id=organization_id,
            ip_address=get_client_ip(request),
            session_id=request.session.session_key
        )
        # validate the scenario doesn't contain vulgar terms using the OpenAPI moderator
        if is_inappropriate(scenario):
            consequence_text = 'Scenario contains inappropriate terms'
            return JsonResponse({'consequence_text': consequence_text}, status=400)

        validation_prompt = f"""
        Evaluate if the following text represents a coherent and plausible cybersecurity scenario, or OT incident scenario, or industrial incident scenario: '{scenario}'. Respond yes if valid, no if invalid.
        """

        validation_response = openai.ChatCompletion.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": validation_prompt}
            ],
            max_tokens=50,
            temperature=0.7
        )

        if "yes" in validation_response['choices'][0]['message']['content'].lower():

            pha_id = request.POST.get('phaID')

            try:
                cyber_pha = tblCyberPHAHeader.objects.get(ID=pha_id)
            except tblCyberPHAHeader.DoesNotExist:
                return JsonResponse({'error': 'PHA record not found'}, status=404)
            facility = Facility.objects.get(id=cyber_pha.facility_id)
            facility_type_id = FacilityType.objects.get(ID=facility.type_id)
            facility_type = facility_type_id.FacilityType
            industry = facility.industry
            chemical_profile = facility.chemicalSummary
            zone = cyber_pha.AssessmentZone
            unit = cyber_pha.AssessmentUnit
            address = cyber_pha.facilityAddress
            country = cyber_pha.country
            devices = cyber_pha.otherSummary

            if malware != '':
                malware_data_query = f"Details about {malware} malware. Ignore irrelevant detail"
                retrieved_chunks = query_index(malware_data_query)
                malware_data = get_summarized_chunks(retrieved_chunks)
            else:
                malware_data = ''
            # Construct a prompt for GPT-4
            system_message = f"""
            You are are expert in cyber-physical scenarios for industrial facilities.  Analyze and consider the following scenario as part of an OT/ICS focused Cyber HAZOPS/Layer of Protection Analysis (LOPA) assessment: {scenario} which occurs at a {facility_type} in the {industry} industry, located at {address} in {country}. 
            The attacker is assumed to be: {attacker}. Malware data is {malware_data}. The attack vector is assumed to be {attack_vector}. The risk category is assumed to be {riskCategory}. Vulnerable systems with Internet exposed IP address {exposed_system}. Vulnerable systems with default or weak credentials {weak_credentials}. Chemicals profile for the facility: {chemical_profile}
            Considering the likely presence of these OT devices: {devices}, concisely describe in 50 words in a list format (separated by semicolons) of a maximum of 5 of the most likely direct consequences of the given scenario. 
            The direct consequences should be specific to the facility and the mentioned details. 
            
            Additional instruction: output ONLY the list items with no text either before or after the list items. DO NOT INCLUDE ANY NOT-TEXT CHARACTERS
            """
            user_message = scenario

            # Query OpenAI API
            response = openai.ChatCompletion.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": system_message},
                    {"role": "user", "content": user_message}
                ],
                max_tokens=100,
                temperature=0.1
            )

            # Extract and process the text from the response
            consequence_text = response['choices'][0]['message']['content']
            consequence_list = consequence_text.split(';')  # Splitting based on the chosen delimiter

            bia_prompt = f""" You are an expert in cyber-physical scenarios for industrial facilities. 
                Given the following consequences from an OT/ICS focused Cyber HAZOPS/Layer of Protection Analysis (LOPA) assessment: 
                Consequences: {consequence_text} Scenario details: {scenario}. 
                Assess the business impact of these consequences under the following 9 categories: safety, danger to life, finance, operations, data, reputation, supply chain, environment, and regulatory compliance. 
                For each category, provide an integer value from 0 to 10 representing the severity of the impact, where 0 means no impact and 10 means maximum impact. 
                Output the result as a JSON object with the categories as keys and integers as values. No additional characters outside of the JSON. """
            bia_response = openai.ChatCompletion.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": bia_prompt},
                    {"role": "user", "content": user_message}
                ],
                max_tokens=100,
                temperature=0.1
            )
            biaData = bia_response['choices'][0]['message']['content']

            attack_tree_json = {}

            if not attack_tree_drawn:
                attack_tree_system_message = """
                  Generate a hierarchical structure of a probable attack tree, based on the MITRE ATT@CK framework for Industrial Control Systems (ICS) applied to and specific to the given OT cybersecurity scenario, in a strictly valid JSON format. 
                  Incorporate relevant terminology from ISA 62443-3-2 if applicable. 
                  The structure should use 'name' for node labels and 'children' for nested nodes, where each node represents a step or method in the attack. 
                  The attack tree must have at least two main branches, each potentially containing dozens of branches or sub-branches. 
                  CRITICAL INSTRUCTION: Ensure the output is in JSON format WITH NO additional characters outside of the JSON structure. The JSON structure should be formatted as: {'name': 'Node Name', 'children': [{...}]}.

                  Example of a correctly formatted output:
                  {
                    "name": "Attack Root",
                    "children": [
                      {
                        "name": "Branch 1",
                        "children": [
                          {
                            "name": "Sub-branch 1.1",
                            "children": []
                          },
                          {
                            "name": "Sub-branch 1.2",
                            "children": []
                          }
                        ]
                      },
                      {
                        "name": "Branch 2",
                        "children": [
                          {
                            "name": "Sub-branch 2.1",
                            "children": []
                          },
                          {
                            "name": "Sub-branch 2.2",
                            "children": []
                          }
                        ]
                      }
                    ]
                  }

                  Please generate a similar structure for the provided cybersecurity scenario, adhering STRICTLY to the JSON format and ensuring at least two main branches are present.
                  """
                user_message = user_message + f'Malware data: {malware_data}'
                # Query OpenAI API for the attack tree
                attack_tree_response = openai.ChatCompletion.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": attack_tree_system_message},
                        {"role": "user", "content": user_message}
                    ],
                    max_tokens=1600,
                    temperature=0.3
                )

                # Process the response for attack tree
                attack_tree_raw = attack_tree_response['choices'][0]['message']['content']
                attack_tree_raw = attack_tree_raw.strip()

                match = re.search(r'\{.*\}', attack_tree_raw, re.DOTALL)
                if match:
                    cleaned_json_str = match.group(0)
                else:
                    cleaned_json_str = "{}"  # Fallback to empty JSON object if no match

                    # Parse the raw JSON string into a Python dictionary
                attack_tree_json = json.loads(cleaned_json_str)

            return JsonResponse({'consequence': consequence_list, 'attack_tree': attack_tree_json, 'biaData': biaData})


    else:
        return JsonResponse({'consequence': [], 'attack_tree': {}, 'biaData': {}, 'error': 'Not a valid scenario'},
                            status=400)

    return JsonResponse({'error': 'Invalid request'}, status=400)


def generate_attack_tree(user_message):
    attack_tree_system_message = f"""
        Generate a hierarchical structure of a potential attack tree for the given cybersecurity scenario in a machine-readable JSON format. The structure should use 'name' for node labels and 'children' for nested nodes. Each node should represent a step or method in the attack, formatted as: {{'name': 'Node Name', 'children': [{{...}}]}}. EXTRA INSTRUCTION: Output MUST be in JSON format with no additional characters outside of the JSON structure.
    """

    # Query OpenAI API for the attack tree
    attack_tree_response = openai.ChatCompletion.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": attack_tree_system_message},
            {"role": "user", "content": user_message}
        ],
        max_tokens=800,
        temperature=0.3
    )

    # Process the response for attack tree
    attack_tree_raw = attack_tree_response['choices'][0]['message']['content']

    try:
        # Parse the raw JSON string into a Python dictionary
        return json.loads(attack_tree_raw)
    except json.JSONDecodeError:
        return {"error": "Invalid JSON format from AI response"}


def is_inappropriate(text):
    openai_api_key = get_api_key('openai')
    openai.api_key = openai_api_key

    try:
        response = openai.Moderation.create(
            input=text,
            model="text-moderation-latest"
        )
        # Check the 'flagged' field in the results
        return any(result['flagged'] for result in response['results'])
    except Exception as e:
        return False  # or handle error appropriately


def write_to_audit(user_id, user_action, user_ip, cyberPHAID=None, cyberPHAScenario=None, qraw=None):
    try:
        user_profile = UserProfile.objects.get(user=user_id)

        audit_log = auditlog(
            user=user_id,
            timestamp=timezone.now(),
            user_action=user_action,
            user_ipaddress=user_ip,
            user_profile=user_profile,
            cyberPHAID=cyberPHAID,
            cyberPHAScenario=cyberPHAScenario,
            qraw=qraw
        )
        audit_log.save()
    except UserProfile.DoesNotExist:
        # Handle the case where UserProfile does not exist for the user
        pass


@csrf_exempt
@require_POST
def assign_cyberpha_to_group(request):
    cyberpha_id = request.POST.get('cyberpha_id')
    existing_group_id = request.POST.get('existing_group_id')
    new_group_name = request.POST.get('new_group_name')
    new_group_type = request.POST.get('new_group_type')
    org_id = get_user_organization_id(request)

    try:
        cyberpha = tblCyberPHAHeader.objects.get(pk=cyberpha_id)

        if existing_group_id:

            group = CyberPHA_Group.objects.get(pk=existing_group_id)
            # Check if the group is already assigned
            if group not in cyberpha.groups.all():
                cyberpha.groups.add(group)
            else:
                return JsonResponse({'status': 'error', 'message': 'CyberPHA is already assigned to this group.'})

            write_to_audit(
                user_id=request.user.id,
                user_ip=get_client_ip(request),
                user_action=f'Add CyberPHA to group: {new_group_name}/{new_group_type}',
                cyberPHAID=cyberpha
            )

        elif new_group_name and new_group_type:
            # Check if group with the same name and type already exists
            organization_instance = Organization.objects.get(id=org_id)
            group, created = CyberPHA_Group.objects.get_or_create(name=new_group_name, group_type=new_group_type,
                                                                  organization=organization_instance)
            if not created:
                return JsonResponse({'status': 'error', 'message': 'Group with this name and type already exists.'})
            cyberpha.groups.add(group)

            write_to_audit(
                user_id=request.user.id,
                user_ip=get_client_ip(request),
                user_action=f'Add CyberPHA to group: {new_group_name}/{new_group_type}',
                cyberPHAID=cyberpha
            )

        return JsonResponse({'status': 'success'})
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)})


@login_required
def fetch_groups(request):
    # Assuming UserProfile contains the organization_id linked to the user
    try:
        user_profile = UserProfile.objects.get(user=request.user)
        organization_id = user_profile.organization_id
    except UserProfile.DoesNotExist:
        # Handle case where user profile or organization is not found
        return JsonResponse({'error': 'User profile or organization not found'}, status=404)

    # Fetch user IDs for the given organization ID
    organization_users = get_organization_users(organization_id)

    cyberpha_id = request.GET.get('cyberpha_id')
    try:
        cyberpha = tblCyberPHAHeader.objects.get(pk=cyberpha_id, UserID__in=organization_users)
    except tblCyberPHAHeader.DoesNotExist:
        # Handle case where tblCyberPHAHeader does not exist or does not belong to the organization
        return JsonResponse({'error': 'CyberPHA not found or does not belong to the organization'}, status=404)

    # Fetch groups associated with this tblCyberPHAHeader that belong to the organization
    groups = cyberpha.groups.all()

    group_data = [{'name': group.name, 'id': group.id} for group in groups]
    return JsonResponse({'groups': group_data})


@login_required
def fetch_all_groups(request):
    # Assuming UserProfile links users to organizations and contains a reference to the user model
    try:
        user_profile = UserProfile.objects.get(user=request.user)
        organization_id = user_profile.organization_id
    except UserProfile.DoesNotExist:
        # Handle case where user profile or organization is not found
        return JsonResponse({'error': 'User profile or organization not found'}, status=404)

    # Assuming tblCyberPHAHeader.UserID is meant to store user ID, but actually stores organization ID
    # and there's a way to map users to their organization IDs correctly in your application
    # This fetches tblCyberPHAHeader records belonging to the user's organization
    # Note: This approach needs adjustment if UserID does not directly relate to organization_id
    organization_user_ids = UserProfile.objects.filter(organization_id=organization_id).values_list('user_id',
                                                                                                    flat=True)
    cyberphas = tblCyberPHAHeader.objects.filter(UserID__in=organization_user_ids).distinct()

    # Fetch groups associated with these tblCyberPHAHeader records
    groups = CyberPHA_Group.objects.filter(tblcyberphaheader__in=cyberphas).distinct()

    group_data = [{'name': group.name, 'id': group.id} for group in groups]
    return JsonResponse({'groups': group_data})


def retrieve_scenario_builder(request, scenario_id):
    try:
        # Retrieve the scenario from the database
        scenario = ScenarioBuilder.objects.get(id=scenario_id)

        # Parse the stored JSON data
        scenario_data = json.loads(scenario.scenario_data)

        # Extract elements from the scenario data
        attack_tree_data = scenario_data.get('attackTree')
        scenario_description = scenario_data.get('scenario')
        investment_projection = scenario_data.get('investment_projection')

        industry = scenario_data.get('industry')
        facility = scenario_data.get('facility')
        country = scenario_data.get('country')
        org = scenario_data.get('org')
        regs = scenario_data.get('regs')
        attacker = scenario_data.get('attacker')
        vector = scenario_data.get('vector')
        target = scenario_data.get('target')
        effect = scenario_data.get('effect')
        network = scenario_data.get('network')
        impact = scenario_data.get('impact')
        motivation = scenario_data.get('motivation')

        # Correctly process consequences to ensure each starts with a single dash
        raw_consequences = scenario_data.get('consequences', '')
        # Split by any known delimiter and ensure each consequence starts with a dash
        consequences_list = raw_consequences.replace(',', '\n').split('\n')
        consequences = [f"- {line.strip()}" if not line.strip().startswith('-') else line.strip()
                        for line in consequences_list if line.strip()]
        # Join the consequences with a line break between each
        formatted_consequences = '\n'.join(consequences)
        table_data = scenario_data.get('tableData')

        costs = scenario_data.get('costs')
        # Check if cost_projection is None or not present, and set a default value if so
        cost_projection = scenario_data.get('cost_projection', '0|0|0|0|0|0|0|0|0|0|0|0')
        if not cost_projection:
            cost_projection = '0|0|0|0|0|0|0|0|0|0|0|0'

        # Prepare factors data
        factors = {}
        for item in table_data:
            factor = item.get('factor')
            score = int(item.get('score').split('/')[0])  # Extract score as integer
            narrative = item.get('narrative')
            factors[factor] = {'score': score, 'narrative': narrative}

        # Return the extracted data
        return JsonResponse({
            'attack_tree_data': attack_tree_data,
            'scenario_description': scenario_description,
            'consequences': formatted_consequences,
            'factors': factors,
            'costs': costs,
            'cost_projection': cost_projection,
            'investment_projection': investment_projection,
            'industry': industry,
            'facility': facility,
            'country': country,
            'org': org,
            'regs': regs,
            'attacker': attacker,
            'vector': vector,
            'target': target,
            'effect': effect,
            'network': network,
            'impact': impact,
            'motivation': motivation
        })

    except ScenarioBuilder.DoesNotExist:
        return JsonResponse({'error': 'Scenario not found'}, status=404)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


logger = logging.getLogger(__name__)


def get_facilities_by_zip(zip_code):
    url = f"https://ofmpub.epa.gov/frs_public2/frs_rest_services.get_facilities?zip_code={zip_code}&output=JSON"
    try:
        response = requests.get(url, verify=False, timeout=20)
        if response.status_code == 200:
            return response.json()
        else:
            logger.error("HTTP Error %s: %s", response.status_code, response.reason)
            return None
    except requests.RequestException as e:
        logger.error("Request error: %s", e)
        return None
    except ValueError as e:
        logger.error("JSON decoding error: %s", e)
        return None


def facilities(request):
    zip_code = request.GET.get('zip')
    if zip_code:
        facilities = get_facilities_by_zip(zip_code)
        if facilities:
            return JsonResponse({'facilities': facilities}, safe=False)
        else:
            return JsonResponse({'error': 'No facilities found or error in API call'}, status=500)
    else:
        return JsonResponse({'error': 'No ZIP code provided'}, status=400)


def get_coordinates_from_address(address, country, google_maps_api_key):
    geocode_url = f"https://maps.googleapis.com/maps/api/geocode/json?address={address},+{country}&key={google_maps_api_key}"
    response = requests.get(geocode_url)
    data = response.json()
    if data['status'] == 'OK':
        latitude = data['results'][0]['geometry']['location']['lat']
        longitude = data['results'][0]['geometry']['location']['lng']
        return latitude, longitude
    else:
        return None, None


def air_quality_index(request):
    latitude = request.GET.get('lat')
    longitude = request.GET.get('lon')
    address = request.GET.get('address')
    country = request.GET.get('country')

    aqicn_api_key = get_api_key("aqicn")  # Replace with your actual AQICN API key

    # If latitude and longitude are not provided, use address and country to get them
    if not latitude or not longitude:
        if address and country:
            latitude, longitude = get_coordinates_from_address(address, country,
                                                               "AIzaSyBJu4p9r_vFL9g5nzctO4yLbNxjK08q4G0")
            if not latitude or not longitude:
                return JsonResponse({'error': 'Failed to geocode address'}, status=400)

    url = f"https://api.waqi.info/feed/geo:{latitude};{longitude}/?token={aqicn_api_key}"

    try:
        response = requests.get(url)
        data = response.json()
        aqi = data.get('data', {}).get('aqi')
        return JsonResponse({'aqi': aqi})
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def facility_risk_profile_newrecord(userid, facility_id, assessment_id, sl, has_ir_plan, ir_plan_never,
                                    ir_plan_tested_date):
    language = 'en'

    facility = get_object_or_404(Facility, id=facility_id)
    industry = facility.industry.Industry
    facility_type = facility.type.FacilityType
    address = facility.address
    shift_model = facility.shift_model
    employees = facility.employees

    if has_ir_plan == 'on':
        ir_plan = 'true'
    else:
        ir_plan = 'false'

    if ir_plan_never == 'on':
        ir_plan_never_tested = 'true'
    else:
        ir_plan_never_tested = 'false'

    if not industry or not facility_type:
        error_msg = "Missing industry or facility type. Complete all fields to get an accurate assessment"
        return JsonResponse({
            'safety_summary': error_msg,
            'chemical_summary': error_msg,
            'physical_security_summary': error_msg,
            'other_summary': error_msg
        })

    openai_api_key = get_api_key('openai')
    openai_model = get_api_key('OpenAI_Model')
    # openai_api_key = os.environ.get('OPENAI_API_KEY')
    openai.api_key = openai_api_key
    context = f"You are an industrial safety and hazard expert. For the {facility} {facility_type} at {address} in the {industry} industry, with {employees} employees working a {shift_model} shift model, (NOTE ALSO - has an OT Cybersecurity Incident Response Plan: {has_ir_plan}. Incident response plan never tested: {ir_plan_never_tested}). "

    prompts = [
        ### f"{context} List safety hazards, max 100 words. - Specific to facility - mechanical or chemical or electrical or heat or cold or crush or height - Space between bullets. \n\nExample Format:\n 1. Specific safety hazard.\n 2. Another specific safety hazard.",
        ### f"{context} List expected chemicals, max 100 words. - Specific to facility - Chemical names only - raw materials and by-products and stored chemicals - Space between bullets. \n\nExample Format:\n 1. Chemical name (raw material or by-product).\n- 2. Another chemical name (raw material or by-product).",
        ### f"{context}, List physical security requirements for the given facility and location - access control - surveillance - consideration of local crime statistics - blind spots - proximity to other infrastructure . Max of 100 words .\n\nExample Format:\n 1. Physical security challenge.\n 2. Another physical security challenge.",
        f"{context}, list of specialized OT and IoT devices and equipment expected to be at the facility. Max of 150 words .\n\nExample Format:\n 1. OT or IoT device (purpose of device).\n 2. Another OT or IoT device (purpose of device).",
        ### f"{context}, list of national and international regulatory compliance containing cybersecurity requirements relevant to the {industry} industry that applies to {facility_type} facilities in {country} . Includes laws and standards. Max of 150 words .\n\nExample Format:\n 1. Compliance name (name of issuing authority).\n 2. Another compliance name (name of issuing authority).",
        ### f"{context}: You are a safety inspector. For a {facility_type} in {country}, estimate a detailed and nuanced safety and hazard risk score. Use a scale from 0 to 100, where 0 indicates an absence of safety hazards and 100 signifies the presence of extreme and imminent fatal hazards. Provide a score reflecting the unique risk factors associated with the facility type and its operational context in {country}. Scores should reflect increments of 10, with each decile corresponding to escalating levels of hazard severity and likelihood of occurrence given the expected attention to safety at the facility. Base your score on a typical {facility_type} in {country}, adhering to expected standard safety protocols, equipment conditions, and operational practices. Provide the score as a single, precise number without additional commentary."
    ]

    responses = []

    # Loop through the prompts and make an API call for each one
    def fetch_response(prompt):
        return openai.ChatCompletion.create(
            # model="gpt-4",
            model=openai_model,
            messages=[
                {"role": "system",
                 "content": "You are a model trained to provide concise responses. Please provide a concise numbered bullet-point list based on the given statement."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1,
            max_tokens=600

        )

        # Use ThreadPoolExecutor to parallelize the API calls

    with concurrent.futures.ThreadPoolExecutor() as executor:
        responses = list(executor.map(fetch_response, prompts))

        # Extract the individual responses
    ### safety_summary = responses[0]['choices'][0]['message']['content'].strip()
    ### chemical_summary = responses[1]['choices'][0]['message']['content'].strip()
    ### physical_security_summary = responses[2]['choices'][0]['message']['content'].strip()
    other_summary = responses[0]['choices'][0]['message']['content'].strip()
    ### compliance_summary = responses[4]['choices'][0]['message']['content'].strip()
    ### pha_score = responses[5]['choices'][0]['message']['content'].strip()

    # Call to facility_threat_profile
    threatSummary, insightSummary, strategySummary = facility_threat_profile(sl, facility, facility_type, address,
                                                                             industry,
                                                                             other_summary,
                                                                             '',
                                                                             has_ir_plan, ir_plan_never_tested,
                                                                             ir_plan_tested_date, '', '', 0)

    return {
        'other_summary': other_summary,
        'threatSummary': threatSummary,
        'insightSummary': insightSummary,
        'strategySummary': strategySummary
    }


@csrf_exempt
def delete_pha_record(request):
    if request.method == 'POST':
        record_id = request.POST.get('id')
        try:
            pha_record = tblCyberPHAHeader.objects.get(ID=record_id)
            pha_record.Deleted = 1
            pha_record.save()

            return JsonResponse({'deleted': True})
        except tblCyberPHAHeader.DoesNotExist:
            return JsonResponse({'deleted': False})
    return JsonResponse({'deleted': False})


def assessment_summary(assessment_id, facilityType, industry):
    openai.api_key = get_api_key('openai')  # Ensure the API key is correctly set

    try:
        assessment = SelfAssessment.objects.get(pk=assessment_id)
    except SelfAssessment.DoesNotExist:
        return "Assessment not found."

    # Preparing the chat messages for the conversation with GPT-4
    messages = []
    pinecone_query = f"Risk assessment, NIST, 62443, C2M2, NIST CSF, gap analysis"
    retrieved_chunks = query_index(pinecone_query)
    summarized_chunks = get_summarized_chunks(retrieved_chunks)
    documents_context = "\n\n".join(summarized_chunks)
    # System message to set the context for the AI
    messages.append({
        "role": "system",
        "content": f"pinecone index data for reference: {documents_context}. You are a cybersecurity analyst and this is a cybersecurity assessment summary and scoring task in the context of a {facilityType} in the {industry} industry. Analyze the provided responses to the cybersecurity assessment questions and write a summary and assign a control effectiveness score out of 100."
    })

    # Adding questions and answers to the conversation
    framework_questions = AssessmentQuestion.objects.filter(framework=assessment.framework).prefetch_related(
        'assessmentanswer_set')
    answered_questions = assessment.answers.all()

    for question in framework_questions:
        answer = answered_questions.filter(question=question).first()
        if answer:
            response_text = "Yes" if answer.response else "No"
            effectiveness = f"Effectiveness: {answer.effectiveness}%" if answer.response and answer.effectiveness is not None else "Effectiveness not applicable."
            message_content = f"Question: {question.text} Answer: {response_text}. {effectiveness}"
        else:
            message_content = f"Question: {question.text} Answer: Unanswered."
        messages.append({"role": "user", "content": message_content})

    # Final user message prompting for the summary and score
    messages.append({
        "role": "user",
        "content": f"Given the responses to the assessment questions, give a concise summary in under 100 words of the cybersecurity program's state in the context of OT Cybersecurity for a {facilityType}, and offer an overall score of control effectiveness out of 100 IMPORTANT: it is vital to consider both the summary and score in the context of the facility and industry. Do not include the facility type in the summary description. i.e. Instead of stating The program at the facilitytype is...  you would instead say The OT Cybersecurity program is ... etc etc . Write the output as two variables in a manner that is easily parsed for display <integer score>|<text summary> "
    })

    # Sending the chat completion request to OpenAI
    try:
        response = openai.ChatCompletion.create(
            model=get_api_key('OpenAI_Model'),
            messages=messages,
            temperature=0.3,
            max_tokens=4096,  # Adjusted for comprehensive analysis
        )
        # Extracting and returning the AI's summary and score
        result = response.choices[0].message['content']

        return result
    except Exception as e:
        return f"Failed to generate assessment due to an API error: {str(e)}"


@require_http_methods(["POST"])
@csrf_exempt  # Consider CSRF protection for production
def get_assessment_summary(request):
    assessment_id = request.POST.get('assessment_id')
    facilityType = request.POST.get('facilityType')
    industry = request.POST.get('industry')
    if not assessment_id:
        return JsonResponse({'error': 'Assessment ID is required'}, status=400)

    result = assessment_summary(assessment_id, facilityType, industry)
    if "Failed to generate" in result:
        return JsonResponse({'error': result}, status=500)

    score, summary = result.split('|', 1)  # Splitting based on the expected format
    return JsonResponse({'score': score, 'summary': summary})


def copy_cyber_pha(request, pha_id):
    if request.method == 'POST':
        original_pha = get_object_or_404(tblCyberPHAHeader, ID=pha_id)

        with transaction.atomic():
            # Get a dictionary of the original PHA excluding many-to-many fields
            original_data = model_to_dict(original_pha, exclude=['id'])
            latitude, longitude = get_coordinates_from_address(
                request.POST.get('facilityAddress') + ',' + request.POST.get('facilityCity') + ', ' + request.POST.get(
                    'facilityCode'), request.POST.get('country'),
                "AIzaSyBJu4p9r_vFL9g5nzctO4yLbNxjK08q4G0")
            update_data = {
                'FacilityName': request.POST.get('FacilityName'),
                'facilityAddress': request.POST.get('facilityAddress'),
                'facilityCity': request.POST.get('facilityCity'),
                'facilityState': request.POST.get('facilityState'),
                'facilityCode': request.POST.get('facilityCode'),
                'country': request.POST.get('country'),
                'facilityLat': latitude,
                'facilityLong': longitude
            }

            original_data.update(update_data)
            original_data.pop('ID', None)  # Ensure the ID is not included

            # Exclude many-to-many fields explicitly before creating the new instance
            m2m_fields = {field.name for field in tblCyberPHAHeader._meta.many_to_many}
            for field in m2m_fields:
                original_data.pop(field, None)

            # Create the new PHA header
            new_pha = tblCyberPHAHeader.objects.create(**original_data)

            # Set many-to-many relationships using .set()
            for field_name in m2m_fields:
                m2m_manager = getattr(original_pha, field_name)
                getattr(new_pha, field_name).set(m2m_manager.all())

            # Copy OneToOne and ForeignKey relationships
            # Assume cybersecurity_defaults is a OneToOneField linked to tblCyberPHAHeader
            if hasattr(original_pha, 'cybersecurity_defaults'):
                defaults_data = model_to_dict(original_pha.cybersecurity_defaults, exclude=['id', 'cyber_pha'])
                CyberPHACybersecurityDefaults.objects.create(cyber_pha=new_pha, **defaults_data)

            # Duplicate related scenarios and their nested related objects
            scenarios = tblCyberPHAScenario.objects.filter(CyberPHA=original_pha)
            for scenario in scenarios:
                scenario_data = model_to_dict(scenario, exclude=['ID', 'CyberPHA', 'userID'])
                scenario_data['CyberPHA'] = new_pha
                if scenario.userID:
                    user_instance = User.objects.get(id=scenario.userID.id)  # Ensure you have the correct user instance
                    scenario_data['userID'] = user_instance
                new_scenario = tblCyberPHAScenario.objects.create(**scenario_data)

                # Correctly duplicating related ScenarioConsequences
                consequences = ScenarioConsequences.objects.filter(scenario=scenario)
                for consequence in consequences:
                    consequence_data = model_to_dict(consequence, exclude=['id', 'scenario'])
                    # Now 'scenario' is not duplicated in consequence_data
                    ScenarioConsequences.objects.create(scenario=new_scenario, **consequence_data)

                    # Correctly duplicating related PHA_Safeguard
                    safeguards = PHA_Safeguard.objects.filter(scenario=scenario)
                    for safeguard in safeguards:
                        safeguard_data = model_to_dict(safeguard, exclude=['id',
                                                                           'scenario'])  # Exclude 'scenario' to avoid conflict
                        PHA_Safeguard.objects.create(scenario=new_scenario, **safeguard_data)

                    # Observations
                    observations = PHA_Observations.objects.filter(scenario=scenario)
                    for observation in observations:
                        observation_data = model_to_dict(observation, exclude=['id', 'scenario'])
                        PHA_Observations.objects.create(scenario=new_scenario, **observation_data)

        return redirect('OTRisk:iotaphamanager')  # Redirect after successful duplication

    return render(request, 'OTRisk/iotaphamanager.html')


def parse_ai_response(ai_response):
    """Parse the AI response into a structured dictionary with score."""
    pattern = re.compile(r'^(\d+)\|([^|]+)\|([^|]+)\|(\d+)$', re.MULTILINE)
    gaps = []
    for match in pattern.finditer(ai_response):
        number, heading, description, score = match.groups()
        cleaned_heading = heading.strip().replace('**', '')
        gaps.append({
            'number': number.strip(),
            'heading': cleaned_heading,
            'description': description.strip(),
            'score': int(score.strip())
        })
    return gaps


@csrf_exempt
def assessment_gap_analysis(request):
    # Extract parameters from POST request
    assessment_id = request.POST.get('assessment_id')
    framework_name = request.POST.get('framework_name')

    if not assessment_id or not framework_name:
        return JsonResponse({'error': 'Missing necessary parameters'}, status=400)

    try:
        assessment = SelfAssessment.objects.get(pk=assessment_id)

    except SelfAssessment.DoesNotExist:
        return JsonResponse({'error': 'Assessment not found'}, status=404)

    messages = [
        {
            "role": "system",
            "content": f"Given the results of a cybersecurity self-assessment for an operational technology environment, interpret these results in light of the '{framework_name}' standards. Consider the objectives, scope, and requirements of the '{framework_name}' and identify where the assessment results indicate potential deficiencies. Describe each deficiency in a structured format: 'number|heading|description|score', with a subjective severity score out of 10 where 0 is 100% lack of alignment and 10 is 100% complete alignment"
        }
    ]
    # Include assessment answers; ideally from a detailed perspective
    assessment_answers = assessment.answers.all()
    for answer in assessment_answers:
        response_text = "Yes" if answer.response else "No"
        effectiveness = f"{answer.effectiveness}%" if answer and answer.effectiveness is not None else "Not applicable"
        messages.append({
            "role": "user",
            "content": f"Question: {answer.question.text}, Answer: {response_text}, Effectiveness: {effectiveness}"
        })

    messages.append({
        "role": "user",
        "content": f"Please summarize the principal gaps between the assessment results and the standards of the mentioned framework. INSTRUCTIONS. The principle sections of '{framework_name}' MUST be used as the heading. For example, in NIST 800-53 ACCESS CONTROL is a principle section and AC-1 Policy and Procedures is a sub-section so ACCESS CONTROL will be heading to use. Description must be 20 words. Include a subjective score for the gap relating to each heading. STRICTLY MAINTAIN THE FORMAT number|heading|description|score"
    })
    openai.api_key = get_api_key('openai')

    try:
        response = openai.ChatCompletion.create(
            model=get_api_key('OpenAI_Model'),
            messages=messages,
            temperature=0.2,
            max_tokens=4000
        )
        ai_response = response.choices[0].message['content']
        gaps = parse_ai_response(ai_response)
        prompt_tokens = response.usage['prompt_tokens']
        completion_tokens = response.usage['completion_tokens']
        total_tokens = response.usage['total_tokens']

        return JsonResponse({'gaps': gaps})  # Send structured data
    except Exception as e:
        return JsonResponse({'error': f"API error: {str(e)}"}, status=500)


@login_required
def load_default_facility(request):
    user_organization_id = get_user_organization_id(request)
    organization_users = get_organization_users(user_organization_id)

    try:

        default_facility = tblCyberPHAHeader.objects.filter(
            UserID__in=organization_users,
            is_default=True,  # Assuming is_default is a BooleanField
            Deleted=0
        ).first()

        investments = CyberSecurityInvestment.objects.filter(cyber_pha_header=default_facility).values(
            'id', 'investment_type', 'vendor_name', 'product_name', 'cost', 'date'
        )
        investments_data = list(investments)
        if not default_facility:
            return JsonResponse({'error': 'No default facility found.'}, status=404)

        # Serialize facility data
        facility_data = {
            'title': default_facility.title,
            'leader': default_facility.PHALeader,
            'leaderemail': default_facility.PHALeaderEmail,
            'unit': default_facility.AssessmentUnit,
            'facility': default_facility.FacilityName,
            'facilitytype': default_facility.FacilityType,
            'assessment_id': default_facility.assessment,
            'zone': default_facility.AssessmentZone,
            'sl_t': default_facility.sl_t,
            'startdate': default_facility.AssessmentStartDate.strftime(
                '%Y-%m-%d') if default_facility.AssessmentStartDate else '',
            'enddate': default_facility.AssessmentEndDate.strftime(
                '%Y-%m-%d') if default_facility.AssessmentEndDate else '',
            'safetysummary': default_facility.safetySummary,
            'chemicalsummary': default_facility.chemicalSummary,
            'physicalsummary': default_facility.physicalSummary,
            'othersummary': default_facility.otherSummary,
            'compliancesummary': default_facility.complianceSummary,
            'threatSummary': default_facility.threatSummary,
            'insightSummary': default_facility.insightSummary,
            'strategySummary': default_facility.strategySummary,
            'country': default_facility.country,
            'shift_model': default_facility.shift_model,
            'EmployeesOnSite': default_facility.EmployeesOnSite,
            'annual_revenue': default_facility.annual_revenue,
            'cyber_insurance': default_facility.cyber_insurance,
            'address': default_facility.facilityAddress,
            'facilityLat': default_facility.facilityLat,
            'facilityLong': default_facility.facilityLong,
            'facilityAQI': default_facility.facilityAQI,
            'facilityCity': default_facility.facilityCity,
            'facilityState': default_facility.facilityState,
            'facilityCode': default_facility.facilityCode,
        }

        return JsonResponse({'headerrecord': facility_data,
                             'investments': investments_data, })

    except Exception as e:

        return JsonResponse({'error': str(e)}, status=500)


def exalens_get_cyberpha_assets(cyberphaid):
    urllib3.disable_warnings(InsecureRequestWarning)

    try:
        cyberpha_header = tblCyberPHAHeader.objects.get(ID=cyberphaid)
        exalens_api_key = cyberpha_header.exalens_api
        exalens_client_id = cyberpha_header.exalens_client
        exalens_ip_address = cyberpha_header.exalens_ip
        assets_url = f"https://{exalens_ip_address}/api/thirdparty/asset"
        headers = {
            'x-client-id': exalens_client_id,
            'x-api-key': exalens_api_key
        }
        assets_response = requests.get(assets_url, headers=headers, verify=False)
        if assets_response.status_code == 200:
            assets_data = assets_response.json().get('data', [])
            if not assets_data:
                return []  # Return an empty list if no data is available
            return assets_data  # Return the list of assets directly
        else:
            return []  # Return an empty list in case of non-200 status codes
    except RequestException as e:
        return []  # Return an empty list if an exception occurs


import urllib3
import datetime
from typing import Any, Dict, Optional
from OTRisk.darktrace_integration.darktrace import Darktrace
from OTRisk.darktrace_integration.device import Device
from OTRisk.darktrace_integration.model_breach import ModelBreach
from OTRisk.darktrace_integration.ai_analyst_incident import AIAnalystIncident
from OTRisk.darktrace_integration.endpoint import Endpoint
from collections import Counter


def darktrace_asset_summary_info(target_asset_ip: str, darktrace_host: str, darktrace_public_token: str,
                                 darktrace_private_token: str) -> Optional[Dict[str, str]]:
    # Disable SSL warnings
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    # Initialize the Darktrace session
    dt_session = Darktrace.initialize_from_string(
        host=darktrace_host,
        public_token=darktrace_public_token,
        private_token=darktrace_private_token,
        timeout=30  # Increase the timeout value
    )

    try:
        devices = Device.get_assets(dt_session)
        asset_detail = next((device for device in devices if device.ip == target_asset_ip), None)

        if asset_detail:
            device_id = asset_detail.did

            # Fetch incidents related to the asset
            incidents = AIAnalystIncident.get_incidents(
                session=dt_session,
                init_date=datetime.datetime.strptime("2023-01-01 00:00:00", "%Y-%m-%d %H:%M:%S"),
                end_date=datetime.datetime.now()
            )

            related_incidents = [incident for incident in incidents if
                                 target_asset_ip in incident.breach_identifiers or
                                 asset_detail.hostname in incident.breach_identifiers]
            critical_incidents = [incident for incident in related_incidents if incident.score >= 90]

            model_breaches = ModelBreach.get_model_breaches(session=dt_session, device=asset_detail)

            total_breach_events = 0
            # Collect detailed model breach information
            model_breaches_details = []
            for breach in model_breaches:
                breach_data = breach.raw["model"]["then"]
                model_breaches_details.append({
                    "model_name": breach_data.get("name", "N/A"),
                    "breach_time": breach.raw.get("breach_time", "N/A"),
                    "score": breach_data.get("score", "N/A"),
                    "description": breach_data.get("description", "N/A"),
                    "actions": breach_data.get("actions", {})
                })
                total_breach_events += 1

            breach_count = len(model_breaches_details)
            if breach_count > 0:
                most_frequent_breach_type = \
                    Counter(breach["model_name"] for breach in model_breaches_details).most_common(1)[0][0]
                highest_score_breach = max(model_breaches_details, key=lambda x: x["score"])
                earliest_breach_time = min(breach["breach_time"] for breach in model_breaches_details)
                latest_breach_time = max(breach["breach_time"] for breach in model_breaches_details)
            else:
                most_frequent_breach_type = "N/A"
                highest_score_breach = {"model_name": "N/A", "score": "N/A", "description": "N/A"}
                earliest_breach_time = "N/A"
                latest_breach_time = "N/A"

            summary_info = {
                "ip_address": asset_detail.ip,
                "asset_type": asset_detail.typename,
                "hostname": asset_detail.hostname,
                "total_events": total_breach_events,
                "breach_count": breach_count,
                "incident_count": len(related_incidents),
                "critical_incident_count": len(critical_incidents),
                "most_frequent_breach_type": most_frequent_breach_type,
                "highest_score_breach": highest_score_breach,
                "earliest_breach_time": earliest_breach_time,
                "latest_breach_time": latest_breach_time
            }

            return summary_info
    except Exception as e:
        logging.error(f"Error fetching asset details: {e}")

    return None


def darktrace_asset_detail(target_asset_ip: str, darktrace_host: str, darktrace_public_token: str,
                           darktrace_private_token: str) -> Optional[Dict[str, str]]:
    # Disable SSL warnings
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    # Initialize the Darktrace session
    dt_session = Darktrace.initialize_from_string(
        host=darktrace_host,
        public_token=darktrace_public_token,
        private_token=darktrace_private_token,
        timeout=30  # Increase the timeout value
    )

    try:
        devices = Device.get_assets(dt_session)
        asset_detail = next((device for device in devices if device.ip == target_asset_ip), None)

        if asset_detail:
            device_id = asset_detail.did

            # Fetch incidents related to the asset
            incidents = AIAnalystIncident.get_incidents(
                session=dt_session,
                init_date=datetime.datetime.strptime("2023-01-01 00:00:00", "%Y-%m-%d %H:%M:%S"),
                end_date=datetime.datetime.now()
            )
            related_incidents = [incident for incident in incidents if
                                 target_asset_ip in incident.breach_identifiers]

            model_breaches = ModelBreach.get_model_breaches(session=dt_session, device=asset_detail)

            # Collect detailed model breach information
            model_breaches_details = []
            for breach in model_breaches:
                breach_data = breach.raw["model"]["then"]
                model_breaches_details.append({
                    "model_name": breach_data.get("name", "N/A"),
                    "breach_time": breach.raw.get("breach_time", "N/A"),
                    "score": breach_data.get("score", "N/A"),
                    "description": breach_data.get("description", "N/A"),
                    "actions": breach_data.get("actions", {})
                })

            breach_count = len(model_breaches_details)
            if breach_count > 0:
                most_frequent_breach_type = \
                    Counter(breach["model_name"] for breach in model_breaches_details).most_common(1)[0][0]
                highest_score_breach = max(model_breaches_details, key=lambda x: x["score"])
                earliest_breach_time = min(breach["breach_time"] for breach in model_breaches_details)
                latest_breach_time = max(breach["breach_time"] for breach in model_breaches_details)
            else:
                most_frequent_breach_type = "N/A"
                highest_score_breach = {"model_name": "N/A", "score": "N/A", "description": "N/A"}
                earliest_breach_time = "N/A"
                latest_breach_time = "N/A"

            summary_info = {
                "ip_address": asset_detail.ip,
                "asset_type": asset_detail.typename,
                "hostname": asset_detail.hostname,
                "breach_count": breach_count,
                "most_frequent_breach_type": most_frequent_breach_type,
                "highest_score_breach": highest_score_breach,
                "earliest_breach_time": earliest_breach_time,
                "latest_breach_time": latest_breach_time
            }

            endpoint_details = Endpoint.get_details(session=dt_session, ip=target_asset_ip)
            largest_data_transfers = Endpoint.get_largest_data_transfers(session=dt_session, ip=target_asset_ip)

            tags = asset_detail.get_tags()

            # Process and format the asset detail for readability
            asset_info = {
                "ip_address": asset_detail.ip,
                "asset_type": asset_detail.typename,
                "hostname": asset_detail.hostname,
                "first_seen": datetime.datetime.fromtimestamp(
                    asset_detail.first_seen / 1000).isoformat() if asset_detail.first_seen else None,
                "last_seen": datetime.datetime.fromtimestamp(
                    asset_detail.last_seen / 1000).isoformat() if asset_detail.last_seen else None,
                "os": asset_detail.os,
                "tags": tags,
                "related_incidents_count": len(related_incidents),
                "related_incidents_types": list(set(incident.summary for incident in related_incidents)),
                "model_breaches_count": len(model_breaches),
                "model_breaches_types": list(set(breach.raw["model"]["then"]["name"] for breach in model_breaches)),
                "endpoint_details_summary": endpoint_details.get("ip", "N/A"),
                "largest_data_transfers_summary": len(largest_data_transfers)
            }

            # Create a condensed summary
            condensed_info = f"""
                IP: {asset_info['ip_address']}
                Type: {asset_info['asset_type']}
                Host: {asset_info['hostname']}
                First Seen: {asset_info['first_seen']}
                Last Seen: {asset_info['last_seen']}
                OS: {asset_info['os']}
                Tags: {', '.join(asset_info['tags'])}
                Incidents: {asset_info['related_incidents_count']} ({', '.join(asset_info['related_incidents_types'])})
                Breaches: {asset_info['model_breaches_count']} ({', '.join(asset_info['model_breaches_types'])})
                Endpoint Details: IP: {asset_info['endpoint_details_summary']}
                Data Transfers: {asset_info['largest_data_transfers_summary']}
            """.strip()

            return {
                "condensed_info": condensed_info,
                "summary_info": summary_info
            }
    except Exception as e:
        logging.error(f"Error fetching asset details: {e}")

    return None


@login_required()
def generate_cyberpha_scenario_description(request):
    if request.method == 'POST':
        # Extracting existing form data
        cyberPHAID = request.POST.get('cyberPHA_ID')

        cyberPHAID = int(cyberPHAID)
        attack_vector = request.POST.get('attackVector', '').strip()
        attacker = request.POST.get('attacker', '').strip()
        target_asset = request.POST.get('targetAsset', '').strip()
        target_asset_purpose = request.POST.get('targetAssetPurpose', '').strip()
        malware = request.POST.get('malware')

        cyberpha_header = get_object_or_404(tblCyberPHAHeader, ID=cyberPHAID)

        facilityid = cyberpha_header.facility_id
        facility = get_object_or_404(Facility, id=facilityid)
        facility_profile = facility.type
        employees_on_site = facility.employees
        address = facility.address
        industry = facility.industry.Industry
        facility_type = facility.type.FacilityType
        chemical_profile = facility_profile.chemical_profile

        connectionFlag = int(request.POST.get('connectionFlag'))

        asset_data = {}
        if connectionFlag == 1:
            darktrace_host = cyberpha_header.darktrace_client
            darktrace_public_token = cyberpha_header.darktrace_public_api
            darktrace_private_token = cyberpha_header.darktrace_private_api

            # Fetch asset details from Darktrace
            asset_data = darktrace_asset_detail(target_asset, darktrace_host, darktrace_public_token,
                                                darktrace_private_token)

        if connectionFlag == 2:
            exalens_client = cyberpha_header.exalens_client
            exalens_api = cyberpha_header.exalens_api
            exalens_url = cyberpha_header.exalens_ip

            # Fetch asset details from Darktrace
            asset_data = exalens_asset_detail(target_asset, exalens_url, exalens_api,
                                              exalens_client)
        network_risk_status = 'unknown'
        if cyberpha_header.darktrace_risk and cyberpha_header.darktrace_risk != '':
            network_risk_status = cyberpha_header.darktrace_risk

        if cyberpha_header.exalens_risk and cyberpha_header.exalens_risk != '':
            network_risk_status = cyberpha_header.exalens_risk

        if malware != '':
            malware_data_query = f"Details about {malware} malware. Ignore irrelevant detail"
            retrieved_chunks = query_index(malware_data_query)
            malware_data = get_summarized_chunks(retrieved_chunks)
        else:
            malware_data = ''

        pinecone_query = f"CyberPHA, Hazops, PHA, CyberHAZOPs, OT Cybersecurity Incidents, OSHA, {industry}"
        retrieved_chunks = query_index(pinecone_query)
        summarized_chunks = get_summarized_chunks(retrieved_chunks)
        documents_context = "\n\n".join(summarized_chunks)

        # Constructing the prompt
        prompt = f"""
        
        You are an OT Cybersecurity expert responsible for safeguarding operations for a {facility_type}. Write a technical OT cybersecurity scenario for a CYBERPHA/ Cyber HAZOPS assessment using LOPA methodology, detailing a credible event specific to manipulation, control, or subversion of given asset. The narrative must be technical, factual, specific to the details given, and concise. Write up to a maximum of 250 words. IMPORTANT: DO NOT describe the consequences, long-term impacts, or make any assumptions about operational disruption and what the scenario means is lacking. No preamble or additional narrative. No repeating of input variables:

        - Attacker: {attacker}
        - Attack Vector: {attack_vector}
        - Country: {address}
        - Industry: {industry}
        - Facility Type: {facility_type}
        - Employee count: {employees_on_site}
        - Affected asset: {target_asset}
        - Chemical profile for the facility: {chemical_profile}
        - Purpose of affected asset: {target_asset_purpose}
        - Asset data reported by threat management: {asset_data}
        - Network work stats: {network_risk_status}
        - Malware data: {malware_data}
        - Pinecone index data: {documents_context}.

        Use the given information above including the data reported by threat management and the information from the pinecone index. Generate a realistic feasible scenario focusing solely on the purpose and type of asset and the potential for that asset to be compromised and how it might be manipulated resulting in a bad outcome in the context of the other given information. INSTRUCTIONS a) If no asset detail is given then focus on the type of facility. b) Do not speculate on mitigation or describe the facility in detail. c) Do not repeat information about the targeted asset or other given detail in the output. d) Do not repeat information that the user is familiar with: industry, number of employees, country, facility type are for context  in generating the scenario and MUST NOT BE REPEATED IN THE RESPONSE. e) Use precise and concise language with as much technical detail as possible given the input details.
        
        The structure of the scenario must be exactly as follows with nothing further:
        <Paragraph describing what is going to happen to the selected asset>.
        <Paragraph describing how the selected asset is affected by the action>.
        
        """

        # Setting OpenAI API key
        openai.api_key = get_api_key('openai')
        open_ai_model = get_api_key('OpenAI_Model')

        # Querying the OpenAI API
        response = openai.ChatCompletion.create(
            model=open_ai_model,
            messages=[
                {"role": "system", "content": prompt}
            ],
            max_tokens=3500,
            temperature=0.1
        )
        # Extracting the generated scenario
        # The response structure is different for chat completions, so adjust accordingly
        if response.choices and response.choices[0].message:
            scenario_description = response.choices[0].message['content'].strip()
        else:
            scenario_description = "No scenario generated."

        return JsonResponse({'scenario_description': scenario_description})
    else:
        return JsonResponse({'error': 'Invalid request'}, status=400)


def summarize_text(text, max_tokens=200):
    response = openai.ChatCompletion.create(
        model='gpt-4o',
        messages=[
            {"role": "system", "content": "Summarize the following text:"},
            {"role": "user", "content": text}
        ],
        max_tokens=max_tokens,
        temperature=0.1
    )
    return response.choices[0].message['content']


def get_summarized_chunks_v2(chunks, max_tokens_per_chunk=200):
    summarized_chunks = []
    for chunk in chunks:
        summarized_text = summarize_text(chunk, max_tokens=max_tokens_per_chunk)
        summarized_chunks.append(summarized_text)
    return summarized_chunks


def get_summarized_chunks(chunks, max_tokens_per_chunk=200):
    summarized_chunks = []
    for chunk in chunks:
        summarized_text = summarize_text(chunk['metadata']['text'], max_tokens=max_tokens_per_chunk)
        summarized_chunks.append(summarized_text)
    return summarized_chunks


import pinecone


def query_index_and_summarize(query, index_name, top_k=7):
    pc = Pinecone(api_key=get_api_key('pinecone'))
    index = pc.Index(index_name)

    # Create an embedding for the query
    response = openai.Embedding.create(input=query, model="text-embedding-ada-002")
    query_embedding = response['data'][0]['embedding']

    # Query Pinecone index
    results = index.query(vector=query_embedding, top_k=top_k, include_metadata=True)

    chunks = [match['metadata']['text'] for match in results['matches']]
    summarized_chunks = get_summarized_chunks_v2(chunks)
    return summarized_chunks


def query_index(query, top_k=7):
    pinecone_api = get_api_key('pinecone')
    pc = Pinecone(api_key=pinecone_api)
    index_name = get_api_key('pinecone_index')
    dimension = 1536  # This should match the dimension of the OpenAI embeddings
    # Create the index if it doesn't exist
    if index_name not in pc.list_indexes().names():
        pc.create_index(
            name=index_name,
            dimension=dimension,
            metric='cosine',
            spec=ServerlessSpec(cloud='aws', region='us-east-1')
        )

    # Connect to the index
    index = pc.Index(index_name)
    # Create an embedding for the query
    response = openai.Embedding.create(input=query, model="text-embedding-ada-002")
    query_embedding = response['data'][0]['embedding']

    # Ensure the query_embedding is a list of floats
    if not isinstance(query_embedding, list) or not all(isinstance(x, float) for x in query_embedding):
        raise ValueError("Query embedding is not in the correct format.")

    # Query Pinecone index
    results = index.query(vector=query_embedding, top_k=top_k, include_metadata=True)

    return results['matches']


def darktrace_assets(cyberphaid):
    # Disable SSL warnings
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    cyberpha_header = tblCyberPHAHeader.objects.get(ID=cyberphaid)
    # Fetch Darktrace credentials from the request or configuration
    darktrace_host = cyberpha_header.darktrace_client
    darktrace_public_token = cyberpha_header.darktrace_public_api
    darktrace_private_token = cyberpha_header.darktrace_private_api

    # Initialize the Darktrace session
    dt_session = DarktraceAPI.initialize_from_string(
        host=darktrace_host,
        public_token=darktrace_public_token,
        private_token=darktrace_private_token
    )

    # Fetch the list of assets
    assets = dt_session.get_assets()

    # Prepare the dataset for the dropdown
    asset_list = []
    for asset in assets:
        asset_data = {
            "ip_address": asset.get("ip"),
            "asset_type": asset.get("typename"),
            "hostname": asset.get("hostname")
        }
        asset_list.append(asset_data)

    # Return the dataset as JSON
    return asset_list


def exalens_asset_detail(target_asset_ip: str, exalens_host: str, exalens_api_key: str, exalens_client_id: str) -> \
        Optional[str]:
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

    incident_url = f"https://{exalens_host}/api/thirdparty/incident/target_ip/{target_asset_ip}?incident=1"
    asset_url = f"https://{exalens_host}/api/thirdparty/asset/ip/{target_asset_ip}"

    headers = {
        'x-client-id': exalens_client_id,
        'x-api-key': exalens_api_key
    }

    try:
        # Fetch incident details
        incident_response = requests.get(incident_url, headers=headers, verify=False)
        if incident_response.status_code != 200:
            raise Exception(f'Failed to fetch incidents, status code: {incident_response.status_code}')

        incident_data = incident_response.json()

        if not incident_data:  # Check if the list is empty
            incident_summary = 'No incidents found for the given IP address.'
        else:
            # Limit to the first three incidents
            limited_incidents = incident_data[:3]
            incident_summaries = []
            for i, incident in enumerate(limited_incidents, start=1):
                summary = (
                    f"Incident {i}\n"
                    f"Incident ID: {incident.get('incident_id', 'N/A')}\n"
                    f"Detection Name: {incident.get('detection_name', 'N/A')}\n"
                    f"Status: {incident.get('status', 'N/A')}\n"
                    f"Severity: {incident.get('severity', 'N/A')}\n"
                    f"Risk Score: {incident.get('risk_score', 'N/A')}\n"
                    f"Summary: {incident.get('summary', 'N/A')}\n"
                    f"Response Recommendations: {incident.get('response_recommendations', 'N/A')}\n"
                )
                incident_summaries.append(summary)

            incident_summary = "\n".join(incident_summaries)

        # Fetch asset details
        asset_response = requests.get(asset_url, headers=headers, verify=False)
        if asset_response.status_code != 200:
            raise Exception(f'Failed to fetch asset details, status code: {asset_response.status_code}')

        asset_data = asset_response.json()

        if not asset_data:  # Check if the list is empty
            asset_summary = 'No asset details found for the given IP address.'
        else:
            asset_info = asset_data[0]  # Assuming the first entry corresponds to the asset
            asset_summary = f"""
                IP: {asset_info.get('ip', 'N/A')}
                VLAN: {asset_info.get('vlan', 'N/A')}
                Status: {asset_info.get('status', 'N/A')}
                OS: {asset_info.get('os', 'N/A')}
                MAC: {asset_info.get('mac', 'N/A')}
                Criticality: {asset_info.get('criticality', 'N/A')}
                Model: {asset_info.get('model', 'N/A')}
                Vendor: {asset_info.get('vendor', 'N/A')}
                Description: {asset_info.get('system_description', 'N/A')}
                Location: {asset_info.get('location', 'N/A')}
                Device Type: {asset_info.get('device_type', 'N/A')}
                First Seen: {datetime.datetime.fromtimestamp(asset_info['first_seen'] / 1000).isoformat() if 'first_seen' in asset_info else 'N/A'}
                Last Seen: {datetime.datetime.fromtimestamp(asset_info['last_seen'] / 1000).isoformat() if 'last_seen' in asset_info else 'N/A'}
            """.strip()

        # Combine the data into a condensed summary
        condensed_info = f"""
            Asset Summary:
            {asset_summary}

            Incident Summary:
            {incident_summary}
        """.strip()

        return condensed_info

    except Exception as e:
        print(f"Error fetching asset details: {e}")
        return None


def create_incident_response_table(steps):
    styles = getSampleStyleSheet()
    helvetica_10 = ParagraphStyle(
        name='Helvetica',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=12,
        alignment=TA_LEFT
    )

    helvetica_bold_10 = ParagraphStyle(
        name='Helvetica-Bold',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        leading=12,
        alignment=TA_LEFT
    )
    table_data = [['Step', 'Title', 'Description']]
    step_counter = 1

    for step in steps:
        if step.strip():
            step = re.sub(r'\*\*|\*\#|\#', '', step).strip()  # Remove any markdown characters
            step_parts = re.split(r'\. ', step, 1)
            if len(step_parts) == 2:
                step_title_description = step_parts[1]
                step_title_parts = step_title_description.split(': ', 1)
                if len(step_title_parts) == 2:
                    step_title = step_title_parts[0]
                    step_description = step_title_parts[1]
                    table_data.append([
                        str(step_counter),
                        Paragraph(step_title, helvetica_bold_10),
                        Paragraph(step_description, helvetica_10)
                    ])
                    step_counter += 1

    col_widths = [0.5 * inch, 2 * inch, 4 * inch]
    incident_response_table = Table(table_data, colWidths=col_widths)

    incident_response_table.setStyle(TableStyle([
        ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('ALIGN', (0, 0), (-1, 0), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('FONTSIZE', (0, 1), (-1, -1), 10),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
    ]))

    return incident_response_table


def add_final_page(elements):
    styles = getSampleStyleSheet()
    helvetica_10 = ParagraphStyle(
        name='Helvetica',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=10,
        leading=12,
        alignment=TA_LEFT
    )

    helvetica_bold_10 = ParagraphStyle(
        name='Helvetica-Bold',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=10,
        leading=12,
        alignment=TA_LEFT
    )
    elements.append(PageBreak())
    final_text = (
        "This playbook is written to address the given scenario for the given facility. The playbook is intended as guidance and a reference, "
        "and is dependent on the various data inputs that have been made to describe the scenario. It is essential that care is taken to verify "
        "and review any actions that are taken within the operational environment based on this playbook to ensure that damage and disruption "
        "is not caused as a consequence of following any of the steps. It is recommended that good change control practices are followed that include "
        "risk assessment, testing, and verification of all changes to network and device configuration before they are deployed into production environments. "
        "The data sources used as reference materials to produce this playbook are updated regularly. If this copy of the playbook is dated more than six months, "
        "it is recommended to log into AnzenOT and print an updated copy."
    )

    elements.append(Spacer(1, 2 * inch))
    elements.append(Paragraph(final_text, helvetica_10))
    return elements


class CustomDocTemplate(BaseDocTemplate):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.page_count = 0
        frame = Frame(self.leftMargin, self.bottomMargin, self.width, self.height, id='normal')
        template = PageTemplate(id='test', frames=frame, onPage=self.add_page_number_and_footer)
        self.addPageTemplates(template)

    def add_page_number_and_footer(self, canvas, doc):
        # Add page number
        page_num_text = f"Page {doc.page}"
        canvas.drawRightString(200 * mm, 10 * mm, page_num_text)
        # Add "Confidential" at the center bottom
        canvas.drawCentredString(self.leftMargin + self.width / 2, 10 * mm, "Confidential")


def generate_scenario_playbook(request):
    if request.method == 'POST':
        data = json.loads(request.body)

        scenario = data.get('scenario')
        recommendations = data.get('recommendations')
        compliancemap = data.get('compliance_map')
        cyberphaid = data.get('cyberphaid')
        scenarioid = data.get('scenarioid')

        # Retrieve the CyberPHAHeader and related Facility data
        cyber_pha_header = get_object_or_404(tblCyberPHAHeader, pk=cyberphaid)
        facility = get_object_or_404(Facility, pk=cyber_pha_header.facility_id)
        facility_type = get_object_or_404(FacilityType, pk=facility.type_id)
        scenario_header = get_object_or_404(tblCyberPHAScenario, pk=scenarioid)
        industry = facility_type.Industry.Industry

        consequences = scenario_header.Consequence

        # Construct detailed queries
        playbook_query = f"incident response playbook, DRAGOS playbook, NIST 800-82"
        risk_assessment_query = f"OT Cybersecurity, Industrial Control Systems"
        incident_response_query = f"incident response procedures, NIST 800-61"
        communication_plan_query = f"communication plan, incident response, comms plans"
        tools_resources_query = f"tools and resources to prevent cybersecurity incidents"

        # Query Pinecone and summarize results
        playbook_data = query_index_and_summarize(playbook_query, "otcyber")
        risk_assessment_data = query_index_and_summarize(risk_assessment_query, "otcyber")
        incident_response_procedures = query_index_and_summarize(incident_response_query, "otcyber")
        communication_plans = query_index_and_summarize(communication_plan_query, "otcyber")
        tools_resources = query_index_and_summarize(tools_resources_query, "otcyber")

        # Prepare the context from Pinecone results
        pinecone_context = {
            "playbook_data": "\n".join(playbook_data),
            "risk_assessment_data": "\n".join(risk_assessment_data),
            "incident_response_procedures": "\n".join(incident_response_procedures),
            "communication_plans": "\n".join(communication_plans),
            "tools_resources": "\n".join(tools_resources),
        }

        # Prepare the prompts for OpenAI API
        messages = {
            "action_items": [
                {"role": "system", "content": "You are an expert in cybersecurity incident response."},
                {"role": "user", "content": (
                    f"""Using the following context from relevant playbooks and risk assessments:\n{pinecone_context['playbook_data']}\n{pinecone_context['risk_assessment_data']}\n
                        Generate a list of up to a maximum of 20 practical, technical, cyber-physical action items for addressing the following scenario: {scenario}. The consequences of the scenario are assumed to be:{consequences}.
                        The action items must be specific, measurable, achievable, relevant, and time-bound (SMART). 
                        Each action item should be written such that an engineer or technology expert can implement it immediately during an ongoing incident without causing further damage or operational loss. 
                        Each action items will include a description of no more than 20 words to support the action item headline
                        Each action items should have a priority - High, Medium, or Low based on the degree of damage and disruption that performing the action item is expected to mitigate. 
                        Do not include action items that require writing new plans, policies, or conducting exercises. 
                        
                        Format the output as follows where , is the delimiter between data fields:
                        <Action Item Number>,<Action Item>,<Action Item Description>,<Priority Level>

                        Example text to illustrate:
                        1,Isolate Compromised PLCs,Immediately disconnect the compromised PLCs from the network to prevent further manipulation of control logic,High
                        2,Activate Backup PLCs: Switch to backup PLCs that have not been compromised to restore control over the mixing and baking processes,Medium

                        Ensure each action item is numbered, followed by a comma, then the action item, followed by a comma, the action item description then comma then the priority level (High, Medium, Low).
                        ADDITIONAL INSTRUCTION: Do not include any markup characters. Do not include any additional narrative, text, or commentary outside of the formatted response."""
                )}
            ],
            "incident_response": [
                {"role": "system", "content": "You are an expert in cybersecurity incident response."},
                {"role": "user", "content": (
                    f"""Based on the following incident response procedures:{pinecone_context['incident_response_procedures']}\n
                            Provide a step-by-step list of up tp a maximum of 20 incident response procedures for detecting, containing, eradicating, and recovering from the scenario: {scenario}. 
                            IMPORTANT: Take into account that the incident is in a {facility_type} in the {industry} industry. 
                            Format the response exactly as described using the guide below and the example as further guidance:

                            <Procedure Heading>
                            <Step number>. <Step Title>: <Step Content>
                            <Step number>. <Step Title>: <Step Content>

                            <Next Procedure Heading>
                            <Step number>. <Step Title>: <Step Content>
                            <Step number>. <Step Title>: <Step Content>

                            And so on. For example:

                            Detection and Analysis
                            1. Incident Detection: Monitor and identify unusual activity.
                            2. Next heading: next description etc etc

                            ADDITIONAL INSTRUCTION: Do not include any markup characters. Do not include any additional narrative, text, or commentary outside of the formatted response."""
                )}
            ],

            "communication_plan": [
                {"role": "system", "content": "You are an expert in cybersecurity compliance."},
                {"role": "user", "content": (
                    f"""Given the following communication plan data: {pinecone_context['communication_plans']}, determine the most effective communication plan for the {facility_type} facility at {facility.address} in the {industry} industry to specifically address the scenario: {scenario}. The estimated consequences of the scenario are : {consequences}. 
                        IMPORTANT: Take into account that the scenario applies to a {facility_type} in the {industry} industry. Take into account the scenario and the consequences in forming the communication plan.

                        Format the response exactly as described using the guide below and the example as further guidance:

                        {{
                            "communication_plan": [
                                {{
                                    "item_number": <Communication plan item number>,
                                    "heading": "<Communication Plan Heading>",
                                    "objective": "<Communication Plan Item Objective>",
                                    "actions": [
                                        {{
                                            "item": "<Action item>",
                                            "description": "<Action description>"
                                        }},
                                        ...
                                    ],
                                    "communication_channels": [
                                        {{
                                            "item": "<Communication Channel>",
                                            "description": "<Channel description>"
                                        }},
                                        ...
                                    ]
                                }},
                                ...
                            ]
                        }}

                        Example for formatting reference. 
                        {{
                            "communication_plan_items": [
                                {{
                                    "item_number": 1,
                                    "heading": "Heading text",
                                    "objective": "Descriptive objective",
                                    "actions": [
                                        {{
                                            "item": "The action to be taken",
                                            "description": "A description of the action to be taken"
                                        }},
                                        {{
                                            "item": "The next action to be taken",
                                            "description": "A description of the next action to be taken"
                                        }},
                                        
                                    ],
                                   and so on for the rest of the communication plan
                                }},
                                ...
                            ]
                        }}

                        ADDITIONAL INSTRUCTION: Do not include any markup characters. Do not include any additional narrative, text, or commentary outside of the formatted response."""
                )}
            ],
        }

        # Query OpenAI API for each prompt
        openai.api_key = get_api_key('openai')

        def query_openai(messages):
            openai.api_key = get_api_key('openai')
            response = openai.ChatCompletion.create(
                model="gpt-4o-mini",
                messages=messages,
                max_tokens=4000,
                temperature=0.2
            )
            return response.choices[0].message['content'].strip()

        action_items = query_openai(messages["action_items"])
        incident_response = query_openai(messages["incident_response"])
        comms_plan = query_openai(messages["communication_plan"])

        # Structure the content for the playbook
        playbook_content = {
            "title": f"CyberPHA Scenario Playbook for {facility.name}",
            "scenario": scenario,
            "recommendations": recommendations,
            "compliance_map": compliancemap,
            "action_items": action_items,
            "incident_response_procedures": incident_response,
            "comms_plan": comms_plan,
            "additional_resources": {
                "playbooks": playbook_data,
                "risk_assessments": risk_assessment_data,
                "incident_response_procedures": incident_response_procedures,
                "communication_plans": communication_plans,
                "tools_resources": tools_resources
            }
        }

        logo_path = os.path.join('static/images', '65C8D0 - Light Blue-2.png')
        buffer = BytesIO()

        doc = CustomDocTemplate(buffer, pagesize=letter)
        elements = []

        styles = getSampleStyleSheet()
        title_style = styles['Title']
        heading_style = styles['Heading1']

        helvetica_10 = ParagraphStyle(
            'Helvetica10',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10,
        )

        # Title page
        elements.append(Spacer(1, 2 * inch))
        elements.append(Paragraph("Scenario Response Playbook", title_style))
        elements.append(Spacer(1, 2 * inch))
        elements.append(Image(logo_path, width=180, height=150))

        elements.append(PageBreak())

        # Scenario page
        elements.append(Paragraph("Scenario", heading_style))
        elements.append(Spacer(1, 0.4 * inch))

        scenario_intro = (
            f"This playbook has been written to address an incident of the following scenario at: {facility.name}, "
            f"a {facility_type.FacilityType} at {facility.address}."
        )
        elements.append(Paragraph(scenario_intro, helvetica_10))
        elements.append(Spacer(1, 0.4 * inch))

        # Add the scenario text inside a box
        scenario_paragraphs = scenario.split('\n\n')
        for paragraph in scenario_paragraphs:
            elements.append(Paragraph(paragraph, helvetica_10))
            elements.append(Spacer(1, 0.2 * inch))

        elements.append(PageBreak())

        # Immediate Actions page
        elements.append(Paragraph("Immediate Actions", heading_style))
        elements.append(Spacer(1, 0.4 * inch))

        action_items_list = playbook_content['action_items'].split('\n')
        table_data = [['#', 'Action Item', 'Description', 'Priority']]

        for rec in action_items_list:
            if rec.strip():
                # Split the action item by its components
                parts = rec.split(',')
                if len(parts) == 4:
                    number = parts[0].strip()
                    action_item = parts[1].strip()
                    description = parts[2].strip()
                    priority = parts[3].strip()

                    table_data.append([
                        number,
                        Paragraph(action_item, helvetica_10),
                        Paragraph(description, helvetica_10),
                        Paragraph(priority, helvetica_10)
                    ])

        # Define the table with specific column widths
        col_widths = [0.05 * doc.width, 0.30 * doc.width, 0.55 * doc.width, 0.10 * doc.width]
        recommendations_table = Table(table_data, colWidths=col_widths)

        # Apply table style
        recommendations_table.setStyle(TableStyle([
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
            ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
            ('ALIGN', (0, 0), (0, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), 6),
            ('RIGHTPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('TOPPADDING', (0, 0), (-1, -1), 12),
        ]))

        elements.append(recommendations_table)
        elements.append(Spacer(1, 0.4 * inch))

        elements.append(PageBreak())
        elements.append(Paragraph("Incident Response Procedures", heading_style))
        elements.append(Spacer(1, 0.4 * inch))

        incident_response_steps = [line.strip() for line in playbook_content['incident_response_procedures'].split('\n')
                                   if line.strip()]
        incident_response_table = create_incident_response_table(incident_response_steps)
        elements.append(incident_response_table)

        elements.append(PageBreak())
        elements.append(Paragraph("Communication Plan", heading_style))
        elements.append(Spacer(1, 0.4 * inch))

        # Correctly load the JSON content
        communication_plan_json = json.loads(playbook_content['comms_plan'])
        if communication_plan_json:
            communication_plan_items = communication_plan_json["communication_plan"]

            table_data = []
            spans = []

            for index, item in enumerate(communication_plan_items):
                item_number = item.get("item_number", "")
                heading = item.get("heading", "")
                objective = item.get("objective", "")
                actions = item.get("actions", [])
                communication_channels = item.get("communication_channels", [])

                # Create bullet points for actions
                actions_bullets = ListFlowable(
                    [ListItem(Paragraph(f"{action['item']}: {action['description']}", helvetica_10)) for action in
                     actions],
                    bulletType='bullet'
                )

                # Create bullet points for communication channels
                channels_bullets = ListFlowable(
                    [ListItem(Paragraph(f"{channel['item']}: {channel['description']}", helvetica_10)) for channel in
                     communication_channels],
                    bulletType='bullet'
                )

                # Add data to the table
                row_start = len(table_data)
                table_data.extend([
                    [str(item_number), Paragraph(heading, helvetica_10), "", ""],
                    ["", Paragraph(objective, helvetica_10), actions_bullets, channels_bullets]
                ])
                row_end = row_start + 1

                # Add span for item number
                spans.append(('SPAN', (0, row_start), (0, row_end)))
                # Add span for heading
                spans.append(('SPAN', (1, row_start), (3, row_start)))

            # Define the column widths for the entire table
            col_widths = [0.05 * doc.width, 0.25 * doc.width, 0.35 * doc.width, 0.35 * doc.width]

            communication_plan_table = Table(table_data, colWidths=col_widths)

            # Apply table style
            table_style = TableStyle([
                ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
                ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                ('ALIGN', (0, 0), (-1, 0), 'LEFT'),
                ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
                ('TOPPADDING', (0, 0), (-1, -1), 12),
            ])

            # Add the dynamic spans to the table style
            for span in spans:
                table_style.add(*span)

            communication_plan_table.setStyle(table_style)

            elements.append(communication_plan_table)
            elements.append(Spacer(1, 0.4 * inch))

        elements = add_final_page(elements)
        # Build the PDF
        doc.build(elements)

        # Get the PDF content
        buffer.seek(0)
        pdf_content = buffer.getvalue()
        buffer.close()

        # Create the HTTP response
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="scenario_playbook.pdf"'
        response.write(pdf_content)

        return response

    return JsonResponse({"error": "Invalid request method."}, status=400)
